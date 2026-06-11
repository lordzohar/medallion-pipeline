"""Append capstone-aligned objective slides to Day_10_Slides.pptx.

Adds one section-divider slide plus one content slide per business
objective requested:
  1. Designing End-to-End Architecture
  2. Integrating Debezium, Kafka, Hop, Airflow
  3. Building a Real-Time Data Pipeline
  4. Deployment and Best Practices
Plus a wrap-up "Capstone in one picture" recap slide.

All content is derived from Capstone_Files/ARCHITECTURE.md and the
docker-compose stack so it matches what bootstrap.ps1 actually deploys.
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DECK = Path(r"C:\Users\Gamer\Documents\GitHub\Medallion pipeline\Week_02\Day_10\Slides\Day_10_Slides.pptx")
prs = Presentation(DECK)

SW, SH = prs.slide_width, prs.slide_height

# Pick a layout: prefer a blank-ish one so we draw exactly what we want.
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=NAVY, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box

def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def slide_header(slide, title, subtitle=None):
    # Top color band
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_text(slide, Inches(0.4), Inches(0.15), SW - Inches(0.8), Inches(0.55),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), SW - Inches(0.8), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

def bullet_list(slide, left, top, width, height, items, *, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = NAVY
    return box

# =============================================================================
# Section divider
# =============================================================================
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SW, SH, NAVY)
add_text(s, Inches(0.5), Inches(2.4), SW - Inches(1.0), Inches(1.2),
         "Capstone — Business Objectives", size=40, bold=True, color=WHITE,
         align=None)
add_text(s, Inches(0.5), Inches(3.6), SW - Inches(1.0), Inches(0.8),
         "End-to-end real-time platform on the Day-10 stack",
         size=18, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text(s, Inches(0.5), Inches(4.5), SW - Inches(1.0), Inches(2.0),
         "1.  Designing End-to-End Architecture\n"
         "2.  Integrating Debezium, Kafka, Hop, Airflow\n"
         "3.  Building a Real-Time Data Pipeline\n"
         "4.  Deployment and Best Practices",
         size=18, color=WHITE)

# =============================================================================
# Objective 1 — Designing End-to-End Architecture
# =============================================================================
s = prs.slides.add_slide(blank_layout)
slide_header(s, "1. Designing End-to-End Architecture",
             "Three real public streams + Postgres CDC → medallion on MinIO → dashboards")

# Left column: design pillars
add_text(s, Inches(0.4), Inches(1.05), Inches(4.6), Inches(0.4),
         "Design pillars", size=15, bold=True, color=ACCENT)
bullet_list(s, Inches(0.4), Inches(1.45), Inches(4.7), Inches(5.5), [
    "Real streams, not simulators — OGN APRS (TCP), NOAA REST, EMSC WebSocket each exercise a different ingestion pattern.",
    "CDC where CDC belongs — Postgres holds only config: regions, alert thresholds, subscriber watchlist. Debezium streams changes so consumers reload without restart.",
    "Medallion in object storage — Avro at bronze + silver (schema-bound), Parquet at gold (analytic-friendly).",
    "Hop is the canonical engine; Python is the executable spec — .hpl / .hwf openable in Hop GUI, Python equivalents run from Airflow so tests don't need a GUI.",
    "Two dashboards, two audiences — quality for DataOps, business for mart consumers.",
    "Everything observable — Prometheus scrapes Kafka JMX, Postgres, MinIO, Connect, both dashboards. Alertmanager routes alerts back into the quality dashboard.",
], size=12)

# Right column: high-level diagram
gx = Inches(5.4); gw = Inches(4.4)
def boxr(top, h, label, fill=LIGHT, color=NAVY, bold=True, size=12):
    add_rect(s, gx, top, gw, h, fill, line=ACCENT)
    add_text(s, gx, top, gw, h, label, size=size, bold=bold, color=color,
             align=None)
    # vertical center hack via padding
add_text(s, gx, Inches(1.05), gw, Inches(0.4),
         "Logical flow", size=15, bold=True, color=ACCENT)

def stage(top, label):
    add_rect(s, gx, top, gw, Inches(0.45), LIGHT, line=ACCENT)
    add_text(s, gx + Inches(0.15), top + Inches(0.08), gw - Inches(0.3),
             Inches(0.3), label, size=12, bold=True)

def arrow(top):
    add_text(s, gx, top, gw, Inches(0.22), "▼", size=12, color=GREY)

y = Inches(1.45)
stages = [
    "Sources: OGN  /  NOAA  /  EMSC  /  config-db",
    "Kafka topics  +  Schema Registry (Avro)",
    "Kafka Connect: Debezium PG  +  2× S3 sink",
    "Bronze (Avro, hourly partitions in MinIO)",
    "Silver (Hop/Python: dedupe, late filter, audit)",
    "Gold (DuckDB → 4 Parquet marts + latest.parquet)",
    "Dashboards: quality / business / live map",
]
for label in stages:
    stage(y, label); y += Inches(0.45)
    if label is not stages[-1]:
        arrow(y); y += Inches(0.22)

# =============================================================================
# Objective 2 — Integrating Debezium, Kafka, Hop, Airflow
# =============================================================================
s = prs.slides.add_slide(blank_layout)
slide_header(s, "2. Integrating Debezium, Kafka, Hop, Airflow",
             "Each tool owns one job; contracts (Avro + topic naming + DAG names) glue them together")

# Four colored cards
cards = [
    ("Debezium",
     ACCENT,
     [
         "Source: config-db (Postgres 15, wal_level=logical) via pgoutput.",
         "Topics: config.public.regions / alert_thresholds / subscriber_watchlist.",
         "Heartbeat topic + DLQ (dlq.config-source) for poison events.",
         "Pinned to debezium-connector-postgres 2.5.4.Final.",
     ]),
    ("Kafka + Connect",
     GREEN,
     [
         "4 stream topics (ogn / noaa.observations / noaa.alerts / seismic) + Debezium CDC topics.",
         "Schema Registry holds Avro subjects, BACKWARD compatibility.",
         "2× S3 sink (one for CDC, one for streams) writes to s3://bronze/<topic>/year=…/hour=…",
         "JMX → :9404 for Prometheus.",
     ]),
    ("Apache Hop",
     RGBColor(0xD6, 0x28, 0x28),
     [
         "Canonical pipelines in .hpl / .hwf (openable in Hop GUI).",
         "bronze_to_silver.py: dedupe on PK, drop events older than max(ts_ms)−24h, write _quality audit.",
         "silver_to_gold.py: DuckDB over silver Avro → 4 Parquet marts (snapshot + latest).",
         "Python reference transforms run from Airflow without the GUI.",
     ]),
    ("Airflow",
     RGBColor(0x8C, 0x56, 0x4B),
     [
         "00_bootstrap (manual)   — create topics, buckets, schemas, connectors.",
         "15_config_drift (2 min)  — mutate config tables to exercise Debezium.",
         "30_hop_medallion (5 min) — silver streams ∥ silver CDC → 4 gold marts.",
         "40_data_quality (5 min) + 50_business_kpis (5 min) — push results to dashboards.",
     ]),
]

card_w = Inches(4.7)
card_h = Inches(2.7)
positions = [
    (Inches(0.4), Inches(1.05)),
    (Inches(5.2), Inches(1.05)),
    (Inches(0.4), Inches(3.95)),
    (Inches(5.2), Inches(3.95)),
]
for (title, color, items), (x, y) in zip(cards, positions):
    add_rect(s, x, y, card_w, Inches(0.42), color)
    add_text(s, x + Inches(0.15), y + Inches(0.06), card_w - Inches(0.3),
             Inches(0.3), title, size=14, bold=True, color=WHITE)
    add_rect(s, x, y + Inches(0.42), card_w, card_h - Inches(0.42), LIGHT,
             line=color)
    bullet_list(s, x + Inches(0.15), y + Inches(0.5), card_w - Inches(0.3),
                card_h - Inches(0.55), items, size=11)

# =============================================================================
# Objective 3 — Building a Real-Time Data Pipeline
# =============================================================================
s = prs.slides.add_slide(blank_layout)
slide_header(s, "3. Building a Real-Time Data Pipeline",
             "From an APRS packet to a refreshed dashboard tile — what runs, where, and how fast")

# Left: the 4 steps with latency budget
steps = [
    ("① Ingest  (continuous services)",
     "OGN / NOAA / EMSC ingestor containers publish Avro records to Kafka.\n"
     "Latency budget: source event → topic  ~  seconds."),
    ("② Land in Bronze  (Kafka Connect)",
     "S3 sink flushes hourly partitions to s3://bronze/<topic>/year=…/hour=….\n"
     "Latency budget: topic → bronze object  ~  ≤ flush interval."),
    ("③ Promote to Silver, then Gold  (Hop / Python, Airflow)",
     "bronze_to_silver: dedupe on PK + keep max(ts_ms), drop late > 24h, write _quality audit.\n"
     "silver_to_gold: DuckDB query → snapshot=YYYYMMDD…/part-0.parquet + latest.parquet.\n"
     "Latency budget: bronze → gold  ~  5 min (30_hop_medallion schedule)."),
    ("④ Serve",
     "quality-dashboard :5001, business-dashboard :5002, live-map :5003.\n"
     "40_data_quality + 50_business_kpis POST to each dashboard every 5 min."),
]
y = Inches(1.05)
for title, body in steps:
    add_rect(s, Inches(0.4), y, Inches(6.4), Inches(0.4), ACCENT)
    add_text(s, Inches(0.55), y + Inches(0.05), Inches(6.2), Inches(0.3),
             title, size=13, bold=True, color=WHITE)
    add_rect(s, Inches(0.4), y + Inches(0.4), Inches(6.4), Inches(0.95),
             LIGHT, line=ACCENT)
    add_text(s, Inches(0.55), y + Inches(0.45), Inches(6.2), Inches(0.9),
             body, size=11)
    y += Inches(1.45)

# Right: the four gold marts the pipeline produces
add_text(s, Inches(7.0), Inches(1.05), Inches(2.9), Inches(0.4),
         "Gold marts produced", size=13, bold=True, color=ACCENT)
marts = [
    ("aircraft_density_by_region",
     "ogn_positions + regions → per-region 1h count, altitude, speed."),
    ("weather_snapshot",
     "noaa_obs + regions joined on station_code → latest per station."),
    ("seismic_24h_summary",
     "seismic_events bucketed by magnitude × region over last 24h."),
    ("region_alert_correlation",
     "Cross-stream: quakes over the live threshold inside any subscribed region."),
]
yy = Inches(1.45)
for name, desc in marts:
    add_rect(s, Inches(7.0), yy, Inches(2.9), Inches(1.2), LIGHT, line=GREY)
    add_text(s, Inches(7.1), yy + Inches(0.05), Inches(2.7), Inches(0.3),
             name, size=12, bold=True, color=GREEN)
    add_text(s, Inches(7.1), yy + Inches(0.4), Inches(2.7), Inches(0.8),
             desc, size=10, color=NAVY)
    yy += Inches(1.30)

# =============================================================================
# Objective 4 — Deployment & Best Practices
# =============================================================================
s = prs.slides.add_slide(blank_layout)
slide_header(s, "4. Deployment & Best Practices",
             "Reproducible startup, observable runtime, recoverable failure modes")

cols = [
    ("Deployment",
     [
         "Single command: bootstrap.ps1 / bootstrap.sh — idempotent, safe to re-run.",
         "Build local image once (day10-app), then `compose pull` so slow remote pulls don't time out `up -d`.",
         "All config from .env; image versions pinned (Kafka 7.5.0, Debezium 2.5.4.Final, Postgres 15).",
         "Plugins downloaded into debezium/plugins (S3 sink + Debezium PG + JMX agent) and mounted into Connect.",
         "Run smoke test: tests/smoke_test.py after bootstrap.",
     ]),
    ("Observability",
     [
         "Prometheus scrapes: Kafka JMX :9404, postgres-exporter :9187, MinIO :9000, Connect REST :8083, both dashboards.",
         "Alert rule files: monitoring/alerts/data_quality_alerts.yml + pipeline_alerts.yml.",
         "Alertmanager → POST quality-dashboard /webhook/alerts so DataOps sees firing alerts in-app.",
         "Grafana auto-provisioned (admin/admin) — pipeline overview dashboard included.",
     ]),
    ("Quality & Recovery",
     [
         "DLQs: dlq.config-source, dlq.s3-sink-bronze-cdc, dlq.s3-sink-bronze-streams.",
         "Silver audit: silver/_quality/<entity>/<iso>.avro on every run.",
         "Schema Registry enforces BACKWARD compatibility — consumers evolve safely.",
         "Debezium resumes from last LSN; S3 sink is idempotent on (topic, partition, offset).",
         "Silver dedupes on PK keeping max(ts_ms); Gold writes snapshot AND latest.parquet so re-runs are non-destructive.",
     ]),
]
xs = [Inches(0.4), Inches(3.7), Inches(7.0)]
for (title, items), x in zip(cols, xs):
    add_rect(s, x, Inches(1.05), Inches(2.9), Inches(0.45), ACCENT)
    add_text(s, x + Inches(0.15), Inches(1.10), Inches(2.7), Inches(0.35),
             title, size=14, bold=True, color=WHITE)
    add_rect(s, x, Inches(1.50), Inches(2.9), Inches(5.4), LIGHT, line=ACCENT)
    bullet_list(s, x + Inches(0.15), Inches(1.6), Inches(2.7), Inches(5.2),
                items, size=10)

# =============================================================================
# Wrap-up: capstone in one picture
# =============================================================================
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Capstone in one picture",
             "What the four objectives produce when you run bootstrap.ps1")

ascii_lines = [
    "OGN APRS    NOAA REST    EMSC WS    config-db (PG, WAL)",
    "    │           │            │              │",
    "    └────────┬──┴─────┬──────┘              │ Debezium pgoutput",
    "             ▼        ▼                     ▼",
    "         ┌─────────────────────────────────────────┐",
    "         │   KAFKA  +  Schema Registry (Avro)      │",
    "         └────────┬──────────────────┬─────────────┘",
    "                  │ S3 sink (streams)│ S3 sink (CDC)",
    "                  ▼                  ▼",
    "         ┌──────────────────────────────┐",
    "         │   MinIO  s3://bronze (Avro)  │",
    "         └──────────────┬───────────────┘",
    "        Hop / Python  bronze_to_silver  (dedupe, audit)",
    "                       ▼",
    "         ┌──────────────────────────────┐",
    "         │   MinIO  s3://silver (Avro)  │",
    "         └──────────────┬───────────────┘",
    "        DuckDB  silver_to_gold  → 4 marts",
    "                       ▼",
    "         ┌──────────────────────────────┐",
    "         │  MinIO  s3://gold  (Parquet) │",
    "         └────┬─────────────────────┬───┘",
    "              ▼                     ▼",
    "   quality-dashboard         business-dashboard",
    "   :5001 (+ /metrics)        :5002 (+ /metrics)",
]
box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), SW - Inches(1.2),
                           SH - Inches(1.4))
tf = box.text_frame
tf.word_wrap = False
for i, line in enumerate(ascii_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = line
    run.font.name = "Consolas"
    run.font.size = Pt(12)
    run.font.color.rgb = NAVY

prs.save(DECK)
print(f"updated: {DECK}")
print(f"total slides now: {len(prs.slides)}")
