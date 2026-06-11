"""Build the full Day-10 capstone training deck from scratch.

Deck structure (52 slides total):

  Opening section          (8 slides)
    1  Title
    2  What you'll learn today
    3  Why this stack — three real-world feeds
    4  The technology stack at a glance (logo grid)
    5  The four modules — your roadmap
    6  Medallion mental model in 30 seconds
    7  How to read this deck
    8  Capstone in one picture

  Module 1  Designing End-to-End Architecture   (10 slides)
  Module 2  Integrating Debezium, Kafka, Hop, Airflow (10 slides)
  Module 3  Building a Real-Time Data Pipeline (10 slides)
  Module 4  Deployment & Best Practices        (10 slides)

  Closing section          (4 slides)
   49  What you built
   50  Where to take it next
   51  Final knowledge check (mixed)
   52  Thank you · references · URL cheat-sheet

Every content slide carries:
  - a real project / tool logo fetched from the web (Wikimedia + project sites)
  - an "Analogy" sidebar box for quick mental hooks
  - bullets sourced from the actual capstone code
  - module-coloured footer bar so you always know where you are

Images are fetched once into _slide_assets/. If a download fails, a labelled
placeholder is drawn instead.
"""
from __future__ import annotations

import io
import os
import sys
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
DECK = ROOT.parent / "Slides" / "Day_10_Slides.pptx"
ASSETS = ROOT / "_slide_assets"
ASSETS.mkdir(exist_ok=True)

# ---- theme ------------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
AMBER  = RGBColor(0xE8, 0x8B, 0x1E)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
PALE   = RGBColor(0xFA, 0xF1, 0xDA)   # analogy box
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ---- image library ----------------------------------------------------------
# Wikimedia thumbnail URLs (stable, SVGs auto-rasterised) plus a couple of
# project-site fallbacks. Keys are short ids referenced by slide specs.
# Wikimedia's Special:FilePath endpoint accepts arbitrary width=&  on any file
# and is the policy-friendly way to pull a thumbnail.
def wm(filename: str, width: int = 400) -> str:
    from urllib.parse import quote
    return f"https://commons.wikimedia.org/wiki/Special:FilePath/{quote(filename)}?width={width}"

# Only niche project logos that appear in this stack. No generic stock photos
# (aircraft / weather / baggage / etc.) — those become hand-drawn analogy
# icons instead, keeping the deck small and offline-friendly.
IMG_URLS: dict[str, str] = {
    "kafka":      wm("Apache kafka.svg", 400),
    "debezium":   "https://debezium.io/assets/images/color_black_debezium_type_600px.png",
    "hop":        "https://hop.apache.org/img/hop-logo.png",
    "airflow":    wm("AirflowLogo.png", 400),
    "postgres":   wm("Postgresql_elephant.svg", 300),
    "docker":     wm("Docker (container engine) logo.svg", 400),
    "grafana":    wm("Grafana_icon.svg", 300),
    "prometheus": wm("Prometheus software logo.svg", 300),
    "minio":      "https://avatars.githubusercontent.com/u/8418256?s=400",
    "duckdb":     wm("DuckDB_logo.svg", 400),
    "avro":       wm("Apache Avro Logo.svg", 300),
    "parquet":    wm("Apache Parquet logo.svg", 400),
}

# Wikimedia requires a UA that identifies the tool AND has a contact handle
# (URL or email). See https://meta.wikimedia.org/wiki/User-Agent_policy
UA = {
    "User-Agent": (
        "DayTenCapstoneSlideBuilder/1.0 "
        "(+https://github.com/local; educational slide deck generator) "
        "python-requests"
    )
}

def fetch(key: str) -> Path | None:
    url = IMG_URLS.get(key)
    if not url:
        return None
    # We don't know the final extension until we see content-type, so use a
    # tentative name and fix it once the response arrives.
    out_png = ASSETS / f"{key}.png"
    out_jpg = ASSETS / f"{key}.jpg"
    if out_png.exists() and out_png.stat().st_size > 0:
        return out_png
    if out_jpg.exists() and out_jpg.stat().st_size > 0:
        return out_jpg
    try:
        r = requests.get(url, headers=UA, timeout=20, allow_redirects=True)
        r.raise_for_status()
        ctype = r.headers.get("content-type", "").lower()
        data = r.content
        if "svg" in ctype or data[:5] == b"<?xml" or data[:4] == b"<svg":
            print(f"[skip] {key}: server returned SVG, can't embed in pptx")
            return None
        # Decide extension from sniffed content
        with Image.open(io.BytesIO(data)) as im:
            fmt = (im.format or "PNG").upper()
        out = out_jpg if fmt in ("JPEG", "JPG") else out_png
        out.write_bytes(data)
        print(f"[ok ] {key} -> {out.name} ({len(data)//1024} KB, {fmt})")
        return out
    except Exception as e:
        print(f"[fail] {key}: {e}")
        for p in (out_png, out_jpg):
            if p.exists():
                p.unlink()
        return None

print("== fetching images ==")
CACHE = {k: fetch(k) for k in IMG_URLS}
print()

# ---- pptx helpers -----------------------------------------------------------
# Fresh deck (built from scratch — the old Northstar slides are dropped).
# Keep the default 10x7.5 slide size that the module-pack layout was tuned for.
prs = Presentation()
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

def add_rect(slide, l, t, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp

def add_text(slide, l, t, w, h, text, *, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, name=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if name:
            r.font.name = name
    return box

def header(slide, module_tag, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.85), NAVY)
    add_text(slide, Inches(0.3), Inches(0.06), Inches(2.2), Inches(0.3),
             module_tag, size=10, bold=True, color=RGBColor(0xCC,0xDD,0xEE))
    add_text(slide, Inches(0.3), Inches(0.28), SW-Inches(0.6), Inches(0.55),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.3), Inches(0.65), SW-Inches(0.6), Inches(0.22),
                 subtitle, size=10, color=RGBColor(0xCC,0xDD,0xEE))

def footer(slide, n_of, module_color):
    add_rect(slide, 0, SH-Inches(0.25), SW, Inches(0.25), module_color)
    add_text(slide, Inches(0.3), SH-Inches(0.23), SW-Inches(0.6), Inches(0.2),
             n_of, size=9, color=WHITE)

# Mapping from "concept" keys used in slide specs to the closest tech logo.
# Lets us keep the original analogy slide content but render real logos.
LOGO_ALIAS = {
    "aircraft":  "kafka",      # OGN slide → Kafka logo
    "weather":   "kafka",      # NOAA / silver slide
    "seismic":   "kafka",
    "baggage":   "docker",
    "warehouse": "minio",
    "factory":   "docker",
}

def resolve_key(key: str) -> str:
    return LOGO_ALIAS.get(key, key)

def insert_image(slide, l, t, w, h, key, fallback_label=""):
    """Drop an image into the box, scaled to fit, or a placeholder."""
    add_rect(slide, l, t, w, h, LIGHT, line=GREY, line_w=Emu(6350))
    real_key = resolve_key(key)
    path = CACHE.get(real_key)
    if path and path.exists():
        try:
            with Image.open(path) as im:
                iw, ih = im.size
            ratio = min((w-Inches(0.15)) / Emu(iw*9525),
                        (h-Inches(0.15)) / Emu(ih*9525))
            tw, th = int(iw*9525*ratio), int(ih*9525*ratio)
            cx = l + (w - tw)//2
            cy = t + (h - th)//2
            slide.shapes.add_picture(str(path), cx, cy, tw, th)
            return
        except Exception as e:
            print(f"[warn] embed {real_key}: {e}")
    add_text(slide, l, t+(h//2)-Inches(0.2), w, Inches(0.4),
             fallback_label or f"[{key}]", size=14, bold=True,
             color=GREY, align=PP_ALIGN.CENTER)

def analogy_box(slide, l, t, w, h, title, body):
    add_rect(slide, l, t, w, h, PALE, line=AMBER, line_w=Emu(12700))
    add_text(slide, l+Inches(0.1), t+Inches(0.08), w-Inches(0.2), Inches(0.32),
             "Analogy", size=10, bold=True, color=AMBER)
    add_text(slide, l+Inches(0.1), t+Inches(0.36), w-Inches(0.2), Inches(0.32),
             title, size=12, bold=True, color=NAVY)
    add_text(slide, l+Inches(0.1), t+Inches(0.7), w-Inches(0.2), h-Inches(0.85),
             body, size=10, color=NAVY)

def bullets(slide, l, t, w, h, items, *, size=12):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size); r.font.color.rgb = NAVY

# ---- generic content slide layout ------------------------------------------
def content_slide(module_tag, color, title, subtitle, img_key, img_label,
                  analogy_title, analogy_body, bullet_items, footer_text):
    s = prs.slides.add_slide(BLANK)
    header(s, module_tag, title, subtitle)

    # Left big content (image on top, bullets below)
    insert_image(s, Inches(0.3), Inches(1.0), Inches(5.6), Inches(3.0),
                 img_key, img_label)
    bullets(s, Inches(0.3), Inches(4.1), Inches(5.6), Inches(1.3),
            bullet_items, size=11)

    # Right side analogy box
    analogy_box(s, Inches(6.1), Inches(1.0), Inches(3.6), Inches(4.4),
                analogy_title, analogy_body)

    footer(s, footer_text, color)
    return s

def diagram_slide(module_tag, color, title, subtitle, draw_fn,
                  analogy_title, analogy_body, footer_text):
    """For slides where the 'image' is a hand-drawn architecture diagram."""
    s = prs.slides.add_slide(BLANK)
    header(s, module_tag, title, subtitle)
    # Diagram canvas on the left
    draw_fn(s, Inches(0.3), Inches(1.0), Inches(5.7), Inches(4.4))
    # Analogy on the right
    analogy_box(s, Inches(6.1), Inches(1.0), Inches(3.6), Inches(4.4),
                analogy_title, analogy_body)
    footer(s, footer_text, color)
    return s

def table_slide(module_tag, color, title, subtitle, headers, rows,
                analogy_title, analogy_body, footer_text):
    s = prs.slides.add_slide(BLANK)
    header(s, module_tag, title, subtitle)
    # Table on the left
    n_cols = len(headers); n_rows = 1 + len(rows)
    tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.0),
                             Inches(5.7), Inches(4.4)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10); r.font.color.rgb = NAVY
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
    analogy_box(s, Inches(6.1), Inches(1.0), Inches(3.6), Inches(4.4),
                analogy_title, analogy_body)
    footer(s, footer_text, color)
    return s

# ---- diagram primitives -----------------------------------------------------
def draw_node(slide, l, t, w, h, label, fill=ACCENT, color=WHITE, size=10):
    add_rect(slide, l, t, w, h, fill, line=NAVY)
    add_text(slide, l, t+(h//2)-Inches(0.13), w, Inches(0.26),
             label, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)

def draw_arrow(slide, x1, y1, x2, y2, label=None):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    line.line.color.rgb = GREY
    line.line.width = Emu(19050)
    if label:
        midx, midy = (x1+x2)//2, (y1+y2)//2
        add_text(slide, midx-Inches(0.6), midy-Inches(0.15),
                 Inches(1.2), Inches(0.3), label, size=8,
                 color=GREY, align=PP_ALIGN.CENTER)

# ---- title-of-module slide --------------------------------------------------
def module_title_slide(module_num, color, title, subtitle, logos):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SH, color)
    add_text(s, Inches(0.5), Inches(0.7), SW-Inches(1.0), Inches(0.5),
             f"MODULE {module_num}", size=14, bold=True,
             color=RGBColor(0xCC,0xDD,0xEE))
    add_text(s, Inches(0.5), Inches(1.2), SW-Inches(1.0), Inches(1.4),
             title, size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(2.6), SW-Inches(1.0), Inches(0.8),
             subtitle, size=16, color=RGBColor(0xCC,0xDD,0xEE))
    # logo strip
    n = len(logos)
    box_w = Inches(1.6); gap = Inches(0.2)
    total = n*box_w + (n-1)*gap
    start = (SW - total)//2
    y = Inches(4.1); h = Inches(1.2)
    for i, key in enumerate(logos):
        x = start + i*(box_w+gap)
        add_rect(s, x, y, box_w, h, WHITE, line=color, line_w=Emu(12700))
        insert_image(s, x+Inches(0.1), y+Inches(0.1),
                     box_w-Inches(0.2), h-Inches(0.2), key, key)
    return s

# ---- quiz slide -------------------------------------------------------------
def quiz_slide(module_tag, color, qa_pairs, footer_text):
    s = prs.slides.add_slide(BLANK)
    header(s, module_tag, "Knowledge Check",
           "5 quick questions — answers in green")

    y = Inches(1.0)
    for i, (q, a) in enumerate(qa_pairs, start=1):
        add_text(s, Inches(0.4), y, SW-Inches(0.8), Inches(0.36),
                 f"Q{i}. {q}", size=12, bold=True, color=NAVY)
        add_text(s, Inches(0.6), y+Inches(0.35), SW-Inches(1.0), Inches(0.42),
                 "A. " + a, size=11, color=GREEN)
        y += Inches(0.82)
    footer(s, footer_text, color)
    return s

# =============================================================================
# Diagram drawers
# =============================================================================
def diagram_medallion(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3), "Medallion in MinIO",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    layer_w = Inches(1.7); layer_h = Inches(0.7); gap = Inches(0.15)
    total = 3*layer_w + 2*gap
    start = l + (w - total)//2
    y = t + Inches(0.45)
    for i, (lbl, fill) in enumerate([
        ("BRONZE\nraw Avro", RGBColor(0xCD,0x7F,0x32)),
        ("SILVER\ndedup Avro", RGBColor(0xC0,0xC0,0xC0)),
        ("GOLD\nParquet marts", RGBColor(0xD4,0xAF,0x37)),
    ]):
        add_rect(s, start+i*(layer_w+gap), y, layer_w, layer_h, fill, line=NAVY)
        add_text(s, start+i*(layer_w+gap), y+Inches(0.08), layer_w, layer_h,
                 lbl, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            ax1 = start+i*(layer_w+gap) + layer_w
            draw_arrow(s, ax1, y+layer_h//2,
                       ax1+gap, y+layer_h//2)
    # legend
    add_text(s, l, t+Inches(1.4), w, Inches(0.3),
             "Promotion: dedup on PK → drop late >24h → DuckDB query → Parquet snapshot + latest",
             size=10, color=GREY, align=PP_ALIGN.CENTER)
    # flow underneath
    add_text(s, l, t+Inches(1.85), w, Inches(0.3),
             "s3://bronze/<topic>/year=…/hour=…/<part>+<off>.avro",
             size=9, color=NAVY, align=PP_ALIGN.CENTER, name="Consolas")
    add_text(s, l, t+Inches(2.15), w, Inches(0.3),
             "s3://silver/<entity>/year=…/day=…/part-<hash>.avro",
             size=9, color=NAVY, align=PP_ALIGN.CENTER, name="Consolas")
    add_text(s, l, t+Inches(2.45), w, Inches(0.3),
             "s3://gold/<mart>/snapshot=…Z/part-0.parquet  +  latest.parquet",
             size=9, color=NAVY, align=PP_ALIGN.CENTER, name="Consolas")
    add_text(s, l, t+Inches(2.95), w, Inches(1.0),
             "Why three layers?\n"
             "Bronze = source-of-truth replay buffer; Silver = clean entity \n"
             "tables (dedup, late-filter, audited); Gold = small, fast marts \n"
             "that dashboards/ML can query directly.",
             size=10, color=NAVY, align=PP_ALIGN.CENTER)

def diagram_cdc(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3), "CDC: how Debezium tails Postgres",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    y = t + Inches(0.5)
    draw_node(s, l+Inches(0.1), y, Inches(1.6), Inches(0.7),
              "Postgres\n(WAL)", fill=RGBColor(0x33,0x67,0x91))
    draw_node(s, l+Inches(2.1), y, Inches(1.6), Inches(0.7),
              "Debezium\npgoutput", fill=ACCENT)
    draw_node(s, l+Inches(4.0), y, Inches(1.6), Inches(0.7),
              "Kafka topic\nconfig.public.*", fill=GREEN)
    draw_arrow(s, l+Inches(1.7), y+Inches(0.35), l+Inches(2.1), y+Inches(0.35),
               "logical")
    draw_arrow(s, l+Inches(3.7), y+Inches(0.35), l+Inches(4.0), y+Inches(0.35),
               "Avro")
    # Detail box
    add_text(s, l, t+Inches(1.5), w, Inches(0.3),
             "Per-row event:", size=11, bold=True, color=NAVY)
    add_text(s, l, t+Inches(1.8), w, Inches(2.4),
             '{ "op": "u",\n'
             '  "before": { "id": 7, "threshold_mag": 4.5 },\n'
             '  "after":  { "id": 7, "threshold_mag": 5.0 },\n'
             '  "source": { "lsn": 0/16D34A8, "ts_ms": 1717... },\n'
             '  "ts_ms":  1717... }\n\n'
             'Order is preserved per-row; resume from last LSN on crash.\n'
             'Snapshot mode bootstraps a new consumer with a full table read.',
             size=10, color=NAVY, name="Consolas")

def diagram_topics(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3), "Topics, partitions, consumers",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # 3 partitions
    y = t + Inches(0.5)
    pw, ph = Inches(1.7), Inches(0.5)
    for i in range(3):
        x = l + Inches(0.1) + i*(pw+Inches(0.15))
        add_rect(s, x, y, pw, ph, ACCENT, line=NAVY)
        add_text(s, x, y+Inches(0.13), pw, Inches(0.3),
                 f"partition {i}", size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # offsets
        for k in range(4):
            ox = x + Inches(0.1) + k*Inches(0.38)
            add_rect(s, ox, y+ph+Inches(0.1), Inches(0.32), Inches(0.32),
                     LIGHT, line=GREY)
            add_text(s, ox, y+ph+Inches(0.12), Inches(0.32), Inches(0.28),
                     str(k), size=9, color=NAVY, align=PP_ALIGN.CENTER)
    # consumer group
    add_text(s, l, t+Inches(1.8), w, Inches(0.3),
             "Consumer group  s3-sink-streams  (3 workers)",
             size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    cy = t + Inches(2.15)
    for i in range(3):
        x = l + Inches(0.1) + i*(pw+Inches(0.15))
        draw_node(s, x, cy, pw, Inches(0.45), f"worker {i}", fill=GREEN)
    # text
    add_text(s, l, t+Inches(2.85), w, Inches(1.5),
             "• A topic is an append-only log split across partitions for parallelism.\n"
             "• A consumer group divides the partitions among its workers; rebalance "
             "happens automatically when a worker joins or dies.\n"
             "• Order is guaranteed inside a partition, never across the whole topic — "
             "pick the partition key carefully (we key by entity id).",
             size=10, color=NAVY)

def diagram_e2e(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3),
             "End-to-end: APRS packet → dashboard tile",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    nodes = [
        ("OGN APRS",     ACCENT),
        ("ingestor",     ACCENT),
        ("Kafka",        GREEN),
        ("S3 sink",      GREEN),
        ("MinIO bronze", AMBER),
        ("Hop silver",   AMBER),
        ("DuckDB gold",  RED),
        ("dashboard",    RED),
    ]
    nw = Inches(1.3); nh = Inches(0.45); ny = t + Inches(0.6)
    total = 4*nw + 3*Inches(0.1)
    start = l + (w-total)//2
    for i, (label, color) in enumerate(nodes[:4]):
        x = start + i*(nw+Inches(0.1))
        draw_node(s, x, ny, nw, nh, label, fill=color, size=10)
        if i < 3:
            draw_arrow(s, x+nw, ny+nh//2, x+nw+Inches(0.1), ny+nh//2)
    ny2 = ny + Inches(0.9)
    for i, (label, color) in enumerate(nodes[4:]):
        x = start + i*(nw+Inches(0.1))
        draw_node(s, x, ny2, nw, nh, label, fill=color, size=10)
        if i < 3:
            draw_arrow(s, x+nw, ny2+nh//2, x+nw+Inches(0.1), ny2+nh//2)
    # join arrow between rows
    draw_arrow(s, start+3*(nw+Inches(0.1))+nw//2, ny+nh,
               start+nw//2, ny2)
    # latency budget
    add_text(s, l, t+Inches(2.3), w, Inches(0.3),
             "Latency budget", size=11, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, l, t+Inches(2.6), w, Inches(1.8),
             "Source → Kafka topic        ~ seconds\n"
             "Topic → bronze object       ≤ S3 sink flush interval\n"
             "Bronze → silver → gold      ~ 5 min   (30_hop_medallion)\n"
             "Gold → dashboard tile       ~ 5 min   (40 + 50 push DAGs)\n\n"
             "End-to-end p95: about 10 minutes from packet to KPI.",
             size=10, color=NAVY, align=PP_ALIGN.CENTER, name="Consolas")

def diagram_compose(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3), "docker-compose stack (22 services)",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # 4 rows of grouped boxes
    groups = [
        ("Streaming", GREEN,
         ["zookeeper","kafka","schema-registry","connect","kafka-ui"]),
        ("Sources",   ACCENT,
         ["config-db","pg-exp","app","ogn","noaa","seismic"]),
        ("Storage",   AMBER,
         ["minio"]),
        ("Compute",   RED,
         ["hop","airflow-db","airflow-init","airflow-web","airflow-sched"]),
        ("Dashboards",ACCENT,
         ["quality","business","live-map"]),
        ("Observ.",   GREY,
         ["prometheus","alertmanager","grafana"]),
    ]
    y = t + Inches(0.5)
    for label, color, items in groups:
        add_rect(s, l, y, Inches(1.0), Inches(0.5), color)
        add_text(s, l, y+Inches(0.12), Inches(1.0), Inches(0.3),
                 label, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        for i, item in enumerate(items):
            x = l + Inches(1.05) + i*Inches(0.78)
            add_rect(s, x, y, Inches(0.75), Inches(0.5), LIGHT, line=color)
            add_text(s, x, y+Inches(0.13), Inches(0.75), Inches(0.3),
                     item, size=8, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER)
        y += Inches(0.6)

def diagram_observability(s, l, t, w, h):
    add_text(s, l, t, w, Inches(0.3), "Observability wiring",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # left: sources
    sources = ["kafka :9404", "pg-exp :9187", "minio :9000",
               "quality :5001", "business :5002", "connect :8083"]
    sy = t + Inches(0.55)
    for s_label in sources:
        add_rect(s, l, sy, Inches(1.6), Inches(0.4), LIGHT, line=ACCENT)
        add_text(s, l, sy+Inches(0.08), Inches(1.6), Inches(0.3),
                 s_label, size=10, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, name="Consolas")
        # arrow to prometheus
        draw_arrow(s, l+Inches(1.6), sy+Inches(0.2),
                   l+Inches(2.4), t+Inches(2.0))
        sy += Inches(0.5)
    # prometheus
    draw_node(s, l+Inches(2.4), t+Inches(1.7), Inches(1.5), Inches(0.7),
              "Prometheus\n:9090", fill=ACCENT)
    # arrow to alertmanager
    draw_arrow(s, l+Inches(3.9), t+Inches(2.05),
               l+Inches(4.4), t+Inches(2.05))
    draw_node(s, l+Inches(4.4), t+Inches(1.7), Inches(1.3), Inches(0.7),
              "Alertmgr\n:9093", fill=RED)
    # arrow back to quality webhook
    draw_arrow(s, l+Inches(5.05), t+Inches(2.4),
               l+Inches(1.6)//2 + l, t+Inches(2.0))
    # grafana
    draw_node(s, l+Inches(2.4), t+Inches(2.8), Inches(1.5), Inches(0.6),
              "Grafana :3000", fill=AMBER)
    draw_arrow(s, l+Inches(2.4)+Inches(0.75), t+Inches(2.4),
               l+Inches(2.4)+Inches(0.75), t+Inches(2.8))
    add_text(s, l, t+Inches(3.6), w, Inches(0.9),
             "Alerts fire from rule_files (data_quality, pipeline) →\n"
             "Alertmanager POSTs JSON to quality-dashboard /webhook/alerts →\n"
             "DataOps sees the alert in the same UI where they see the metric.",
             size=10, color=NAVY, align=PP_ALIGN.CENTER)

# =============================================================================
# Opening section  (8 slides) — replaces the old Northstar starter slides
# =============================================================================

def opening_title_slide():
    s = prs.slides.add_slide(BLANK)
    # full-bleed navy with a coloured side bar
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SH, ACCENT)

    add_text(s, Inches(0.5), Inches(0.55), SW - Inches(1.0), Inches(0.5),
             "DAY 10  ·  CAPSTONE", size=14, bold=True,
             color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, Inches(0.5), Inches(1.05), SW - Inches(1.0), Inches(1.4),
             "Unified Streaming Data Platform",
             size=34, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(2.4), SW - Inches(1.0), Inches(0.9),
             "Real public feeds  →  Kafka  →  Medallion on MinIO  →  Dashboards",
             size=16, color=RGBColor(0xCC, 0xDD, 0xEE))

    # 8-logo strip showcasing the stack
    logos = ["debezium", "kafka", "hop", "airflow",
             "minio", "duckdb", "prometheus", "grafana"]
    n = len(logos)
    box_w, gap, h = Inches(1.05), Inches(0.1), Inches(1.0)
    total_w = n * box_w + (n - 1) * gap
    start_x = (SW - total_w) // 2
    y = Inches(4.3)
    for i, key in enumerate(logos):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y, box_w, h, WHITE, line=ACCENT, line_w=Emu(9525))
        insert_image(s, x + Inches(0.06), y + Inches(0.06),
                     box_w - Inches(0.12), h - Inches(0.12), key, key)

    # bottom band with subtitle
    add_rect(s, 0, SH - Inches(1.4), SW, Inches(1.4), RGBColor(0x07, 0x1C, 0x33))
    add_text(s, Inches(0.5), SH - Inches(1.2), SW - Inches(1.0), Inches(0.4),
             "4 modules  ·  40 lessons  ·  20 knowledge-check questions  ·  one runnable stack",
             size=14, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(0.5), SH - Inches(0.75), SW - Inches(1.0), Inches(0.4),
             "OGN aircraft   ·   NOAA weather   ·   EMSC seismic   ·   Postgres config CDC",
             size=12, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)
    return s


def opening_objectives_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "What you will learn today",
           "Five outcomes, one stack — by the end of Day 10 you can…")
    items = [
        ("Architect a streaming + medallion data platform",
         "Pick the right tool at each layer (CDC vs polling, Avro vs Parquet, "
         "stream vs batch) and defend the choice."),
        ("Integrate Debezium, Kafka, Hop and Airflow",
         "Wire CDC into a broker, push to object storage, refine through "
         "bronze → silver → gold, schedule everything from Airflow."),
        ("Build a real-time pipeline against live public feeds",
         "OGN APRS, NOAA, EMSC — three different ingestion patterns "
         "(TCP, REST, WebSocket) running side-by-side."),
        ("Deploy reproducibly with Docker Compose",
         "One bootstrap script that pins versions, builds local images, "
         "pulls remote ones, waits for health, and seeds metadata."),
        ("Operate with Prometheus / Grafana / Alertmanager",
         "Metrics, alerts, dashboards and DLQs — observability baked in, "
         "not bolted on."),
    ]
    y = Inches(1.05)
    for n, (title, body) in enumerate(items, start=1):
        add_rect(s, Inches(0.3), y, Inches(0.55), Inches(0.55), ACCENT)
        add_text(s, Inches(0.3), y + Inches(0.13), Inches(0.55), Inches(0.3),
                 str(n), size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.0), y, SW - Inches(1.3), Inches(0.32),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(1.0), y + Inches(0.32), SW - Inches(1.3), Inches(0.6),
                 body, size=11, color=GREY)
        y += Inches(0.95)
    footer(s, "Slide 2 / 52  ·  Opening", ACCENT)
    return s


def opening_why_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "Why this stack — three real feeds",
           "Real data is the syllabus; the tools just react to it")

    # left: 3 feed cards
    feeds = [
        ("OGN APRS",  ACCENT, "ogn.aircraft.positions",
         "Bursty TCP stream of glider beacons. Forces us to handle "
         "out-of-order events, partial packets, idle gaps."),
        ("NOAA",      GREEN,  "noaa.observations / noaa.alerts",
         "REST poll every 60 s for station obs + active warnings. "
         "Classic pull-with-cursor pattern."),
        ("EMSC",      AMBER,  "seismic.events",
         "Long-lived WebSocket pushing earthquakes worldwide. Forces "
         "reconnect-with-backoff + de-duplication."),
    ]
    y = Inches(1.05)
    for name, color, topic, body in feeds:
        add_rect(s, Inches(0.3), y, Inches(0.18), Inches(1.05), color)
        add_rect(s, Inches(0.48), y, Inches(5.5), Inches(1.05), LIGHT,
                 line=color, line_w=Emu(9525))
        add_text(s, Inches(0.6), y + Inches(0.08), Inches(5.3), Inches(0.3),
                 name, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.6), y + Inches(0.35), Inches(5.3), Inches(0.28),
                 topic, size=10, color=ACCENT, name="Consolas")
        add_text(s, Inches(0.6), y + Inches(0.6), Inches(5.3), Inches(0.45),
                 body, size=10, color=GREY)
        y += Inches(1.15)

    # right: analogy
    analogy_box(s, Inches(6.1), Inches(1.05), Inches(3.6), Inches(4.35),
                "Why three feeds, not one",
                "Three feeds force three different ingestion patterns into the "
                "same pipeline. A platform that handles TCP, REST AND WebSocket "
                "won't break when feed #4 arrives. Building for one source is "
                "easy; building for variety is the lesson.")

    footer(s, "Slide 3 / 52  ·  Opening", ACCENT)
    return s


def opening_stack_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "The technology stack at a glance",
           "12 tools, each doing one job well")

    cells = [
        ("debezium",   "Debezium",   "CDC source"),
        ("kafka",      "Kafka",      "event broker"),
        ("avro",       "Avro",       "wire schema"),
        ("hop",        "Apache Hop", "ETL pipelines"),
        ("airflow",    "Airflow",    "orchestration"),
        ("postgres",   "Postgres",   "config + WAL"),
        ("minio",      "MinIO",      "S3 object store"),
        ("parquet",    "Parquet",    "gold-mart format"),
        ("duckdb",     "DuckDB",     "in-process SQL"),
        ("docker",     "Docker",     "runtime"),
        ("prometheus", "Prometheus", "metrics"),
        ("grafana",    "Grafana",    "dashboards"),
    ]
    cols, rows = 4, 3
    cell_w = Inches(2.25); cell_h = Inches(1.35)
    grid_w = cols * cell_w + (cols - 1) * Inches(0.15)
    grid_h = rows * cell_h + (rows - 1) * Inches(0.15)
    start_x = (SW - grid_w) // 2
    start_y = Inches(1.15)
    for i, (key, name, role) in enumerate(cells):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + Inches(0.15))
        y = start_y + r * (cell_h + Inches(0.15))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line=ACCENT, line_w=Emu(9525))
        insert_image(s, x + Inches(0.1), y + Inches(0.1),
                     cell_w - Inches(0.2), Inches(0.75), key, name)
        add_text(s, x, y + Inches(0.9), cell_w, Inches(0.25),
                 name, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.1), cell_w, Inches(0.22),
                 role, size=9, color=GREY, align=PP_ALIGN.CENTER)

    footer(s, "Slide 4 / 52  ·  Opening", ACCENT)
    return s


def opening_roadmap_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "The four modules — your roadmap",
           "Each module = 10 slides = one slice of the platform")

    modules = [
        ("MODULE 1", ACCENT, "Architecture",
         "Patterns: medallion, CDC, Kappa.\nWhy real streams, not simulators."),
        ("MODULE 2", GREEN,  "Integration",
         "Debezium · Kafka · Hop · Airflow.\nContracts: Avro + topic names + DAG ids."),
        ("MODULE 3", AMBER,  "Pipeline build",
         "Bronze → Silver → Gold.\nDedupe, late-events, DuckDB marts."),
        ("MODULE 4", RED,    "Deployment & Ops",
         "Compose, Prometheus, Alertmanager.\nGotchas, DLQs, capacity rules."),
    ]
    card_w = Inches(2.25); card_h = Inches(3.3); gap = Inches(0.15)
    total_w = 4 * card_w + 3 * gap
    start_x = (SW - total_w) // 2
    y = Inches(1.2)
    for i, (tag, color, title, body) in enumerate(modules):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, y, card_w, Inches(0.55), color)
        add_text(s, x, y + Inches(0.12), card_w, Inches(0.4),
                 tag, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, x, y + Inches(0.55), card_w, card_h - Inches(0.55),
                 LIGHT, line=color, line_w=Emu(9525))
        add_text(s, x + Inches(0.15), y + Inches(0.7),
                 card_w - Inches(0.3), Inches(0.45),
                 title, size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.2),
                 card_w - Inches(0.3), card_h - Inches(1.35),
                 body, size=10, color=GREY, align=PP_ALIGN.CENTER)
        # 10-slide ribbon
        ry = y + card_h - Inches(0.5)
        add_text(s, x, ry, card_w, Inches(0.3),
                 "10 lessons + quiz", size=10, bold=True, color=color,
                 align=PP_ALIGN.CENTER)

    # bottom: numbered timeline
    ty = Inches(5.7)
    line_y = ty + Inches(0.2)
    line = s.shapes.add_connector(1, Inches(0.6), line_y,
                                  SW - Inches(0.6), line_y)
    line.line.color.rgb = GREY
    line.line.width = Emu(19050)
    stops = [("Start", Inches(0.6)),
             ("M1",    Inches(2.8)),
             ("M2",    Inches(5.0)),
             ("M3",    Inches(7.2)),
             ("M4",    SW - Inches(0.6))]
    for lbl, x in stops:
        add_rect(s, x - Inches(0.1), line_y - Inches(0.1),
                 Inches(0.2), Inches(0.2), ACCENT)
        add_text(s, x - Inches(0.5), line_y + Inches(0.15),
                 Inches(1.0), Inches(0.3),
                 lbl, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    footer(s, "Slide 5 / 52  ·  Opening", ACCENT)
    return s


def opening_medallion_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "Medallion in 30 seconds",
           "Three bands of trust — the mental model you carry into every module")
    # reuse the medallion diagram
    diagram_medallion(s, Inches(0.3), Inches(1.05), Inches(5.7), Inches(4.4))
    analogy_box(s, Inches(6.1), Inches(1.05), Inches(3.6), Inches(4.4),
                "Restaurant kitchen",
                "Bronze = raw ingredients off the truck. Silver = washed, "
                "chopped, mise-en-place. Gold = the plated dish you serve. "
                "Each layer is owned by a different role and is cheap to "
                "redo from the layer below.")
    footer(s, "Slide 6 / 52  ·  Opening", ACCENT)
    return s


def opening_how_to_read_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "How to read every slide in this deck",
           "Same template across all 40 lesson slides — learn it once")

    # Mock content panel on the left
    add_rect(s, Inches(0.3), Inches(1.05), Inches(5.7), Inches(4.4),
             LIGHT, line=GREY, line_w=Emu(6350))
    add_rect(s, Inches(0.3), Inches(1.05), Inches(5.7), Inches(0.55), NAVY)
    add_text(s, Inches(0.45), Inches(1.18), Inches(5.4), Inches(0.4),
             "Lesson title goes here", size=14, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(1.85), Inches(5.3), Inches(2.0),
             WHITE, line=ACCENT, line_w=Emu(9525))
    add_text(s, Inches(0.5), Inches(2.75), Inches(5.3), Inches(0.3),
             "[ niche tool logo ]", size=12, color=GREY,
             align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(0.5), Inches(4.0), Inches(5.3), Inches(1.3),
             "• Bullet sourced from the actual capstone code\n"
             "• Real path, real config, real version\n"
             "• Plus the WHY, not just the WHAT\n"
             "• Module-coloured footer = you-are-here",
             size=11, color=NAVY)

    # Right side: analogy box explanation
    analogy_box(s, Inches(6.1), Inches(1.05), Inches(3.6), Inches(4.4),
                "The Analogy box",
                "Every lesson has one — a real-world parallel that locks the "
                "tech idea into memory. Read it FIRST when a topic feels "
                "abstract; come back to it when teaching someone else.")

    footer(s, "Slide 7 / 52  ·  Opening", ACCENT)
    return s


def opening_one_picture_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "OPENING", "The whole capstone in one picture",
           "Bookmark this slide — every later lesson is a zoom into one box")
    diagram_e2e(s, Inches(0.3), Inches(1.05), Inches(9.4), Inches(4.7))
    footer(s, "Slide 8 / 52  ·  Opening", ACCENT)
    return s


# Build the opening section
opening_title_slide()
opening_objectives_slide()
opening_why_slide()
opening_stack_slide()
opening_roadmap_slide()
opening_medallion_slide()
opening_how_to_read_slide()
opening_one_picture_slide()

# =============================================================================
# Module 1  Designing End-to-End Architecture
# =============================================================================
M1 = "MODULE 1  · Architecture"
M1C = ACCENT
mod = 1; total = 10

module_title_slide(
    1, ACCENT,
    "Designing End-to-End Architecture",
    "Real streams → Kafka → Medallion on MinIO → Dashboards",
    ["kafka","debezium","airflow","minio"],
)

content_slide(
    M1, M1C,
    "Why real streams, not simulators",
    "OGN aircraft, NOAA weather, EMSC seismic — three live public feeds",
    "aircraft", "OGN glider in flight",
    "Stethoscope vs simulator",
    "A medical student learns more from one real heartbeat than from a thousand "
    "perfectly-shaped synthetic ones. Real feeds bring real surprises — gaps, "
    "out-of-order events, schema wiggles — that simulators never do.",
    [
        "OGN APRS over TCP → ogn.aircraft.positions  (bursty, lossy)",
        "NOAA api.weather.gov REST poll → noaa.observations / noaa.alerts",
        "EMSC seismic WebSocket → seismic.events  (long-lived socket)",
        "Each ingestor exercises a different ingestion pattern — TCP / REST / WS.",
    ],
    f"Slide 2 / {total}  ·  Module 1",
)

diagram_slide(
    M1, M1C,
    "The medallion pattern",
    "Bronze / Silver / Gold — bands of trust, not folders",
    diagram_medallion,
    "Restaurant kitchen",
    "Bronze = raw ingredients off the truck. Silver = washed, chopped, prepped "
    "(mise en place). Gold = the plated dish the customer sees. Each layer is "
    "owned by a different role and is cheap to redo from the layer below.",
    f"Slide 3 / {total}  ·  Module 1",
)

content_slide(
    M1, M1C,
    "CDC vs polling for reference data",
    "Why Debezium beats SELECT * every minute",
    "postgres", "PostgreSQL",
    "Doorbell vs door check",
    "Polling is walking to the door every minute to see if anyone is there. "
    "CDC is the doorbell — the database tells you the moment something happens, "
    "and you sleep through the silence.",
    [
        "Postgres holds config tables only (regions, thresholds, watchlist).",
        "wal_level=logical + Debezium pgoutput → row-level change events.",
        "Topics: config.public.regions, .alert_thresholds, .subscriber_watchlist",
        "Consumers reload thresholds without restart — drift DAG (every 2 min) "
        "mutates rows to keep CDC busy in the demo.",
    ],
    f"Slide 4 / {total}  ·  Module 1",
)

content_slide(
    M1, M1C,
    "Schema Registry as the contract",
    "Avro subjects + BACKWARD compatibility = no 3am surprises",
    "avro", "Apache Avro",
    "Building blueprints",
    "Architects don't ship a new building plan without comparing it to the old "
    "one — you can add windows, but you can't move load-bearing walls. Schema "
    "Registry enforces that, automatically, on every produce.",
    [
        "Each topic value has a subject:  <topic>-value",
        "BACKWARD compatibility = new schema can read old data.",
        "Producers serialize against the registered schema or fail loudly.",
        "Bronze keeps Avro (small + schema-bound); Gold flips to Parquet "
        "(column-pruning friendly for dashboards).",
    ],
    f"Slide 5 / {total}  ·  Module 1",
)

table_slide(
    M1, M1C,
    "Topic + partition inventory",
    "Naming convention: <source>.<entity>(.action)?",
    ["Topic", "Parts", "Source", "DLQ"],
    [
        ("ogn.aircraft.positions",            6, "OGN ingestor",     "—"),
        ("noaa.observations",                 3, "NOAA ingestor",    "—"),
        ("noaa.alerts",                       3, "NOAA ingestor",    "—"),
        ("seismic.events",                    3, "EMSC ingestor",    "—"),
        ("config.public.regions",             1, "Debezium CDC",     "dlq.config-source"),
        ("config.public.alert_thresholds",    1, "Debezium CDC",     "dlq.config-source"),
        ("config.public.subscriber_watchlist",1, "Debezium CDC",     "dlq.config-source"),
        ("config.heartbeat",                  1, "Debezium",         "—"),
    ],
    "Highway lanes",
    "Partitions are lanes on a highway. More lanes = more cars in parallel, "
    "but order is only guaranteed within one lane. Pick the partition key "
    "(entity id) so things that must stay in order share a lane.",
    f"Slide 6 / {total}  ·  Module 1",
)

content_slide(
    M1, M1C,
    "Storage layout in MinIO",
    "Hourly partitioned Avro at bronze, snapshot+latest Parquet at gold",
    "minio", "MinIO (S3-compatible)",
    "Library shelves",
    "Bronze is the returns cart — books dumped in arrival order. Silver is the "
    "Dewey-decimal shelf — one copy per ISBN. Gold is the librarian's pick-of-"
    "the-week display — small, curated, swap-out-the-poster fast.",
    [
        "s3://bronze/<topic>/year=YYYY/month=MM/day=DD/hour=HH/...avro",
        "s3://silver/<entity>/year=YYYY/month=MM/day=DD/part-<hash>.avro",
        "s3://silver/_quality/<entity>/<utc-iso>.avro  (audit trail)",
        "s3://gold/<mart>/snapshot=…Z/part-0.parquet  +  latest.parquet",
    ],
    f"Slide 7 / {total}  ·  Module 1",
)

table_slide(
    M1, M1C,
    "Four gold marts, four questions",
    "Each mart is one decision the business wants to make in <1 sec",
    ["Mart", "Inputs", "Grain"],
    [
        ("aircraft_density_by_region",   "ogn + regions",                   "region · last 1h"),
        ("weather_snapshot",             "noaa_obs + regions",              "station_code · latest"),
        ("seismic_24h_summary",          "seismic_events",                  "mag-bucket × region · 24h"),
        ("region_alert_correlation",     "seismic + regions + thresholds + watchlist", "alert event"),
    ],
    "Newsroom front page",
    "A mart is a front-page headline: one number, one chart, ready to print. "
    "The investigative reporting (Avro joins, dedupe) happens upstream so "
    "the editor (dashboard) doesn't think.",
    f"Slide 8 / {total}  ·  Module 1",
)

diagram_slide(
    M1, M1C,
    "Trade-off: kappa over lambda",
    "One pipeline for batch and stream — replay from bronze when needed",
    lambda s,l,t,w,h: (
        add_text(s, l, t, w, Inches(0.3), "Kappa: re-process from the log",
                 size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER),
        draw_node(s, l+Inches(0.5), t+Inches(0.6), Inches(1.6), Inches(0.5),
                  "Kafka log", fill=GREEN),
        draw_node(s, l+Inches(2.7), t+Inches(0.6), Inches(1.6), Inches(0.5),
                  "Stream job v1", fill=ACCENT),
        draw_node(s, l+Inches(4.7), t+Inches(0.6), Inches(0.9), Inches(0.5),
                  "Gold v1", fill=AMBER),
        draw_arrow(s, l+Inches(2.1), t+Inches(0.85), l+Inches(2.7), t+Inches(0.85)),
        draw_arrow(s, l+Inches(4.3), t+Inches(0.85), l+Inches(4.7), t+Inches(0.85)),
        draw_node(s, l+Inches(2.7), t+Inches(1.4), Inches(1.6), Inches(0.5),
                  "Stream job v2", fill=ACCENT),
        draw_node(s, l+Inches(4.7), t+Inches(1.4), Inches(0.9), Inches(0.5),
                  "Gold v2", fill=AMBER),
        draw_arrow(s, l+Inches(2.1), t+Inches(0.85), l+Inches(2.7), t+Inches(1.65)),
        draw_arrow(s, l+Inches(4.3), t+Inches(1.65), l+Inches(4.7), t+Inches(1.65)),
        add_text(s, l, t+Inches(2.2), w, Inches(2.2),
                 "Lambda needs two code paths (batch + stream) that must agree.\n"
                 "Kappa keeps one path; rewrites = a new consumer that re-reads\n"
                 "bronze. Cheaper to test, easier to reason about.\n\n"
                 "In this stack:  the Hop/Python jobs read bronze each run, so a\n"
                 "logic fix is just  ‘delete silver+gold, re-run’.",
                 size=11, color=NAVY, align=PP_ALIGN.CENTER, name="Consolas"),
    ),
    "Tape vs notebook",
    "A tape recorder lets you rewind and re-listen. A notebook only has what "
    "you wrote down. Kappa says: keep the tape (Kafka log) and you can always "
    "re-create the notes.",
    f"Slide 9 / {total}  ·  Module 1",
)

quiz_slide(
    M1, M1C,
    [
        ("In this stack, why does Postgres only hold config tables and not facts?",
         "Facts are high-volume and append-only — they belong in Kafka. Postgres "
         "is for small reference data that CDC streams cheaply."),
        ("Which Avro compatibility mode lets new schemas read old data?",
         "BACKWARD — the default we configure on Schema Registry."),
        ("Why pick Avro for bronze/silver and Parquet for gold?",
         "Avro is row-oriented + schema-bound (good for streaming writes). "
         "Parquet is column-oriented (good for dashboard reads)."),
        ("What does the partition key control?",
         "Which partition an event lands in — and therefore the unit of ordering. "
         "Same key → same partition → guaranteed order."),
        ("If a Gold mart logic bug is found, what's the recovery path?",
         "Delete that mart, re-run silver_to_gold — bronze is the source of "
         "truth and re-derivation is cheap (Kappa)."),
    ],
    f"Slide 10 / {total}  ·  Module 1 — Knowledge Check",
)

# =============================================================================
# Module 2  Integrating Debezium, Kafka, Hop, Airflow
# =============================================================================
M2 = "MODULE 2  · Integration"
M2C = GREEN

module_title_slide(
    2, GREEN,
    "Integrating Debezium, Kafka, Hop, Airflow",
    "Four tools, four jobs, one contract (Avro + topic names + DAG ids)",
    ["debezium","kafka","hop","airflow"],
)

diagram_slide(
    M2, M2C,
    "Debezium deep dive",
    "Tailing the WAL, not querying the table",
    diagram_cdc,
    "Train conductor's logbook",
    "A conductor doesn't ask each passenger who they are — they read the "
    "boarding manifest. The WAL is the manifest of every database change, in "
    "order, and Debezium just reads it.",
    f"Slide 2 / {total}  ·  Module 2",
)

diagram_slide(
    M2, M2C,
    "Kafka topics, partitions, consumer groups",
    "Append-only logs split for parallelism",
    diagram_topics,
    "Supermarket checkouts",
    "One topic = one queue. Partitions = open lanes. Consumer group = cashiers "
    "who agree to split the lanes between them. Add a cashier → automatic "
    "re-balance, no customer left waiting.",
    f"Slide 3 / {total}  ·  Module 2",
)

table_slide(
    M2, M2C,
    "Schema Registry compatibility modes",
    "Pick once per subject — drives what producers can change",
    ["Mode", "New reader sees old?", "Old reader sees new?", "Use when"],
    [
        ("BACKWARD",     "yes (default)", "no",          "consumers upgrade after producers"),
        ("FORWARD",      "no",            "yes",         "producers upgrade after consumers"),
        ("FULL",         "yes",           "yes",         "both must work, tightest"),
        ("NONE",         "n/a",           "n/a",         "dev only — break-glass"),
    ],
    "Phone-charger USB-C",
    "BACKWARD compatibility is like a new charger that still fits old phones — "
    "the upgrade path nobody has to think about.",
    f"Slide 4 / {total}  ·  Module 2",
)

content_slide(
    M2, M2C,
    "Kafka Connect — source and sink plugins",
    "The bridge between the broker and everything that is not Kafka",
    "kafka", "Apache Kafka Connect",
    "Power adapters",
    "Connectors are travel adapters — same Kafka 'plug', different shapes on "
    "the other side: Postgres, S3, ElasticSearch, Snowflake. You don't write "
    "the adapter; you configure it with JSON.",
    [
        "debezium-connector-postgres  — source: WAL → Kafka",
        "kafka-connect-s3 (×2)        — sink: Kafka → MinIO bronze",
        "Plugins drop into  debezium/plugins/  → mounted at /etc/kafka-connect/plugins.",
        "JSON config posted to REST :8083; restart-safe (offsets in internal topics).",
    ],
    f"Slide 5 / {total}  ·  Module 2",
)

content_slide(
    M2, M2C,
    "Apache Hop — pipelines and workflows",
    "Visual ETL with .hpl / .hwf files that are also git-friendly XML",
    "hop", "Apache Hop",
    "LEGO instruction booklet",
    "A Hop pipeline (.hpl) is the LEGO instruction sheet — every step is a "
    "named brick, ordered, with inputs and outputs you can see. The Python "
    "twin runs the same recipe from Airflow when there's no GUI.",
    [
        "bronze_to_silver.hpl  — per entity dedupe + late-filter + write audit.",
        "silver_to_gold.hpl    — DuckDB query over silver Avro → Parquet.",
        "Workflows (.hwf) chain pipelines: wf_bronze_to_silver, wf_full_medallion.",
        "Python equivalents (hop/transforms/*.py) run headless inside Airflow.",
    ],
    f"Slide 6 / {total}  ·  Module 2",
)

table_slide(
    M2, M2C,
    "Airflow DAGs in this stack",
    "5 DAGs — 1 manual bootstrap, 1 drift generator, 3 every-5-minute jobs",
    ["DAG", "Schedule", "What it does"],
    [
        ("00_bootstrap",       "manual",   "create topics, buckets, schemas, connectors"),
        ("15_config_drift",    "every 2 m","mutate Postgres rows to keep CDC busy"),
        ("30_hop_medallion",   "every 5 m","silver streams ∥ silver CDC → 4 gold marts"),
        ("40_data_quality",    "every 5 m","rule pack → POST quality-dashboard"),
        ("50_business_kpis",   "every 5 m","POST business-dashboard /api/refresh"),
    ],
    "Symphony conductor",
    "Airflow doesn't play any instrument — it tells each section when to come "
    "in and how loud. The streams are still playing all the time on their own.",
    f"Slide 7 / {total}  ·  Module 2",
)

diagram_slide(
    M2, M2C,
    "How they wire together",
    "One picture: every box you saw in the previous slides, in place",
    diagram_e2e,
    "Postal system",
    "Postman picks up letters (ingestor), sorting office routes by zipcode "
    "(Kafka partition), regional warehouse stores them (bronze), couriers "
    "re-package by route (silver→gold), recipient opens the envelope (dashboard).",
    f"Slide 8 / {total}  ·  Module 2",
)

content_slide(
    M2, M2C,
    "Top 5 integration pitfalls (and how this stack dodges them)",
    "Where teams typically lose a week — and what to do instead",
    "docker", "Docker (the runtime that holds it together)",
    "Climbing rope rules",
    "Climbers say 'always tie in, always check, always communicate.' Same here: "
    "always pin versions, always check Schema Registry, always send to a DLQ "
    "rather than crash.",
    [
        "1. Unpinned image tags — :latest changes silently. We pin Kafka 7.5.0, Debezium 2.5.4.Final.",
        "2. No DLQs — poison message kills the consumer. We ship 3 DLQ topics.",
        "3. WAL not enabled — Debezium silently does nothing. 03_enable_wal.sql.",
        "4. Schema drift without registry — bronze breaks silver. BACKWARD compat.",
        "5. DAG owns the stream — when scheduler dies, data dies. Ingestors are services.",
    ],
    f"Slide 9 / {total}  ·  Module 2",
)

quiz_slide(
    M2, M2C,
    [
        ("What does Debezium read from Postgres to produce events?",
         "The write-ahead log (WAL) via the pgoutput logical-decoding plugin — "
         "it does NOT poll the tables."),
        ("Where do Kafka Connect plugin JARs live in this stack?",
         "debezium/plugins/ on the host, mounted into the connect container at "
         "/etc/kafka-connect/plugins (declared via CONNECT_PLUGIN_PATH)."),
        ("Why are the 3 ingestors services and not Airflow DAGs?",
         "Streams are continuous; if Airflow's scheduler dies, you'd lose data. "
         "Long-lived sockets belong in services that restart on failure."),
        ("What's the trade-off of FULL schema compatibility vs BACKWARD?",
         "FULL = old AND new readers both work = safest but tightest constraints. "
         "BACKWARD = only new readers need to read old data = easier to evolve."),
        ("If a poison message lands on a topic, what stops Connect from crashing?",
         "errors.tolerance + the DLQ topic (dlq.s3-sink-bronze-cdc etc.) — the "
         "bad record is parked, processing continues."),
    ],
    f"Slide 10 / {total}  ·  Module 2 — Knowledge Check",
)

# =============================================================================
# Module 3  Building a Real-Time Data Pipeline
# =============================================================================
M3 = "MODULE 3  · Pipeline build"
M3C = AMBER

module_title_slide(
    3, AMBER,
    "Building a Real-Time Data Pipeline",
    "From a public packet on the internet to a refreshed KPI in ~10 minutes",
    ["kafka","minio","duckdb","grafana"],
)

content_slide(
    M3, M3C,
    "OGN ingestor walkthrough",
    "APRS over TCP → Avro on ogn.aircraft.positions",
    "aircraft", "OGN — Open Glider Network",
    "Police radio scanner",
    "An OGN beacon is a tiny radio announcement: 'I'm here, this altitude, this "
    "speed.' The ingestor is the scanner that listens, transcribes, and pins "
    "each call to a noticeboard (Kafka).",
    [
        "ogn-client connects to APRS-IS with a callsign and a radius filter.",
        "Parser handles malformed packets — bad ones go to logs, not topic.",
        "Avro serializer registers schema on first message.",
        "Restart policy: unless-stopped — survives a single APRS server hiccup.",
    ],
    f"Slide 2 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "NOAA + EMSC ingestors",
    "REST polling and a long-lived WebSocket — two more ingestion shapes",
    "weather", "NOAA — National Weather Service",
    "Two ways to read the news",
    "Polling NOAA = refreshing the homepage every 60 s. The seismic WebSocket = "
    "subscribing to a push notification. Same goal (latest news), different cost.",
    [
        "noaa_ingestor: api.weather.gov for stations + active alerts, 60 s poll.",
        "seismic_ingestor: EMSC WebSocket; reconnects with backoff on close.",
        "Both emit to Avro topics, sharing the same Schema Registry.",
        "All 3 read CITY_LAT / CITY_LON from .env so swapping demo location is trivial.",
    ],
    f"Slide 3 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "Bronze: S3 sink partitioning",
    "Why hourly partitions, and what the path actually looks like",
    "parquet", "Object-store path layout",
    "Filing cabinet",
    "Year/month/day/hour folders are the labels on a filing cabinet drawer. "
    "Without them, a query has to read every file. With them, the planner "
    "skips 99% of the cabinet.",
    [
        "Path: s3://bronze/<topic>/year=YYYY/month=MM/day=DD/hour=HH/<topic>+<part>+<off>.avro",
        "Hive-style partitions = automatic pruning in DuckDB / Spark / Athena.",
        "Hour is a sweet spot — small enough to prune, big enough to avoid "
        "millions of tiny files.",
        "Two sink connectors split CDC vs streams to keep tunings independent.",
    ],
    f"Slide 4 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "Silver: dedupe and late events",
    "What 'clean' actually means in code — bronze_to_silver.py",
    "weather", "Silver = clean entity tables",
    "Photographer's keeper folder",
    "A photographer shoots 200 frames of the same pose. Silver is the keeper "
    "folder — one best frame per pose, blurries deleted, EXIF audited.",
    [
        "Dedup: group by natural PK, keep the row with max(ts_ms).",
        "Late filter: drop rows older than max(ts_ms) − 24h (configurable).",
        "Audit: every run writes silver/_quality/<entity>/<utc-iso>.avro "
        "(rows_in, rows_out, dropped_late, dup_collapsed).",
        "Idempotent — re-running just re-writes the same files, no double-count.",
    ],
    f"Slide 5 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "Gold: DuckDB marts",
    "Embedded analytic engine doing the heavy lifting, no extra server",
    "duckdb", "DuckDB — in-process analytics",
    "Pocket calculator vs mainframe",
    "DuckDB is a high-end pocket calculator: no server, no admin, but it can "
    "crunch tens of GB of Parquet over the network in seconds. Perfect for "
    "Gold-mart builds.",
    [
        "silver_to_gold.py reads silver Avro via DuckDB (fastavro fallback).",
        "Joins regions / thresholds / watchlist inside one SELECT.",
        "Writes snapshot=<utc>/part-0.parquet AND overwrites latest.parquet.",
        "Dashboards point at latest.parquet — atomic rename = zero-downtime swap.",
    ],
    f"Slide 6 / {total}  ·  Module 3",
)

diagram_slide(
    M3, M3C,
    "End-to-end latency budget",
    "Where the seconds and minutes actually go",
    diagram_e2e,
    "Pizza delivery time",
    "Order placed → kitchen acknowledges (Kafka commit, seconds) → cooked "
    "(Hop, 5 min) → driver leaves (gold publish) → you eat (dashboard refresh). "
    "10-minute p95 is fine for KPIs, painful for live ops.",
    f"Slide 7 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "Three dashboards, three audiences",
    "Quality, business, and live map — one stack, three viewpoints",
    "grafana", "Grafana / Flask dashboards",
    "Cockpit, weather radar, in-flight movie",
    "Pilot watches the cockpit (quality dashboard). Crew checks weather radar "
    "(business KPIs). Passengers watch the live map (live_map_dashboard). Same "
    "plane, same data, different needs.",
    [
        "quality-dashboard :5001 — rule pass/fail counters, DLQ depth, alerts.",
        "business-dashboard :5002 — KPIs from gold marts, refreshes every 5 min.",
        "live-map :5003 — Leaflet map of OGN aircraft, near-real-time.",
        "All expose /metrics — Prometheus scrapes for unified observability.",
    ],
    f"Slide 8 / {total}  ·  Module 3",
)

content_slide(
    M3, M3C,
    "Big-picture analogy — airport baggage handling",
    "Why the whole architecture clicks once you see it as luggage",
    "baggage", "Baggage belt at Geneva Airport",
    "Tag → belt → x-ray → cart → carousel",
    "Check-in tags the bag (schema). Conveyor sorts by destination (partition "
    "key). X-ray flags problems (DLQ). Cart batches to the plane (bronze "
    "hourly files). Carousel delivers to you (dashboard). Each handler does "
    "ONE thing well — that's the whole game.",
    [
        "Check-in agent     = Schema Registry  (tag the bag once, trust it everywhere)",
        "Sorter / belt      = Kafka partitions",
        "X-ray operator     = Connect with DLQ tolerance",
        "Aircraft hold      = S3 bronze, sealed for transport",
        "Baggage carousel   = Gold mart latest.parquet",
    ],
    f"Slide 9 / {total}  ·  Module 3",
)

quiz_slide(
    M3, M3C,
    [
        ("Why does silver keep an audit file under _quality/?",
         "So you can prove every Gold number — counts, drops, dedupes — without "
         "re-running. Compliance + debugging in one cheap file."),
        ("What rule decides which row wins when silver dedupes?",
         "Group by natural PK; keep the row with the maximum ts_ms (most recent)."),
        ("Why does Gold write BOTH a timestamped snapshot AND latest.parquet?",
         "Snapshot = history / audit. latest.parquet = stable URL the dashboard "
         "reads. Atomic rename means dashboards never see a half-written file."),
        ("If a NOAA ingestor crashes, what is the data loss exposure?",
         "Anything in flight that wasn't yet produced to Kafka. The container "
         "restarts under restart: unless-stopped, then resumes polling."),
        ("What's the end-to-end p95 latency target for this stack?",
         "About 10 minutes — driven by the 5-min Hop schedule plus the 5-min "
         "dashboard-push DAG. Bronze landing itself is seconds."),
    ],
    f"Slide 10 / {total}  ·  Module 3 — Knowledge Check",
)

# =============================================================================
# Module 4  Deployment & Best Practices
# =============================================================================
M4 = "MODULE 4  · Deployment & Ops"
M4C = RED

module_title_slide(
    4, RED,
    "Deployment and Best Practices",
    "Reproducible startup, observable runtime, recoverable failure",
    ["docker","prometheus","grafana","postgres"],
)

diagram_slide(
    M4, M4C,
    "The compose stack at a glance",
    "22 containers, one bridge network, one .env",
    diagram_compose,
    "Apartment block",
    "Each container is an apartment. Compose is the building, the .env is the "
    "rulebook every apartment shares, the bridge network is the hallway. You "
    "move out one tenant without disturbing the others.",
    f"Slide 2 / {total}  ·  Module 4",
)

content_slide(
    M4, M4C,
    "bootstrap.ps1 walkthrough",
    "What `./bootstrap.ps1` actually does on first run vs re-run",
    "docker", "Docker Compose orchestration",
    "Recipe with checkpoints",
    "Like a recipe that says 'if dough is already proven, skip step 4.' "
    "Idempotent bootstrap = safe to re-run after any failure; only the missing "
    "steps execute.",
    [
        "1. Ensure .env exists (copy from .env.example if not).",
        "2. Download Connect plugins + JMX agent (skip if present).",
        "3. Build local app-base image, then compose pull, then compose up -d.",
        "4. Wait-Http on minio / schema-registry / connect / airflow.",
        "5. Create MinIO buckets, register Avro schemas, register all connectors.",
        "6. Unpause the 4 continuous DAGs and print the URL cheat-sheet.",
    ],
    f"Slide 3 / {total}  ·  Module 4",
)

content_slide(
    M4, M4C,
    "Image strategy: pin, build, pull — in that order",
    "Why the 'pull access denied for day10-app' warning is harmless",
    "docker", "Docker image lifecycle",
    "Library card vs your own bookshelf",
    "Pulled images are library books (versioned, returnable). The day10-app "
    "image is a book you wrote yourself — Docker Hub doesn't have it, the "
    "warning is just the library looking and finding nothing.",
    [
        "Pin every external image (Kafka 7.5.0, Debezium 2.5.4.Final, Postgres 15).",
        "`compose build app-base` first — locally-built image now in the cache.",
        "`compose pull --ignore-pull-failures` to bring remote ones without timing out `up -d`.",
        "`compose up -d` — no network races, no surprise version drift.",
    ],
    f"Slide 4 / {total}  ·  Module 4",
)

diagram_slide(
    M4, M4C,
    "Observability: Prometheus + Grafana + Alertmanager",
    "Scrape the metrics, alert on the rules, show the dashboards",
    diagram_observability,
    "Hospital monitor",
    "Sensors on every patient (exporters) → ECG screen (Prometheus) → red "
    "button when vitals cross thresholds (Alertmanager) → nurse dashboard "
    "(quality-dashboard / Grafana).",
    f"Slide 5 / {total}  ·  Module 4",
)

content_slide(
    M4, M4C,
    "DLQ and recovery patterns",
    "Three DLQs in this stack — what fills them, what drains them",
    "kafka", "Kafka — DLQ topics",
    "Hospital triage",
    "DLQs are the triage corner: don't block the ER, treat the easy patients "
    "first, come back to the hard cases with a specialist. Same for messages — "
    "park the poison, keep throughput, fix later.",
    [
        "dlq.config-source             — Debezium can't decode this WAL row",
        "dlq.s3-sink-bronze-cdc        — sink couldn't write this record",
        "dlq.s3-sink-bronze-streams    — sink couldn't write this record",
        "Drain: re-process by hand → republish to original topic OR archive forever.",
        "Gauge: quality_dlq_size{topic} alerts when growth > 0.",
    ],
    f"Slide 6 / {total}  ·  Module 4",
)

table_slide(
    M4, M4C,
    "Common production gotchas",
    "Failure mode → first symptom → quickest fix",
    ["Failure mode", "First symptom", "Quickest fix"],
    [
        ("WAL not enabled on Postgres",      "Debezium creates slot then 0 events",       "ALTER SYSTEM wal_level = logical; restart"),
        ("Schema Registry down",             "Connect produces ‘SerializationException’", "compose restart schema-registry; replay DLQ"),
        ("Connect plugin path wrong",        "0 connectors visible at /connectors",       "check CONNECT_PLUGIN_PATH and bind mount"),
        ("MinIO out of space",               "S3 sink RUNNING → FAILED, 507 in logs",     "mc rm or extend volume; restart connector"),
        ("Airflow scheduler dead",           "DAGs queued but never start",               "compose restart airflow-scheduler"),
        ("Pinned image moved on Hub",        "compose pull says ‘not found’",             "switch tag, audit version in .env"),
    ],
    "Pilot's checklist",
    "Pilots don't memorise emergencies — they read the laminated card. This "
    "table is your laminated card.",
    f"Slide 7 / {total}  ·  Module 4",
)

content_slide(
    M4, M4C,
    "Secrets, configuration, and the .env file",
    "Where credentials live (and where they must NOT live)",
    "docker", ".env — single source of config truth",
    "Hotel room safe",
    "Your passwords are jewellery — they live in the safe (.env / a secret "
    "manager), not on the dresser (git, Dockerfile, source code). Everything "
    "else (image versions, ports, paths) is fine in the open.",
    [
        ".env holds: MinIO creds, Postgres creds, Airflow Fernet key, ports.",
        ".env.example is committed; .env is .gitignored.",
        "For real prod: swap .env for Docker secrets / Vault / AWS SSM.",
        "Never bake credentials into Dockerfile RUN lines or git history.",
    ],
    f"Slide 8 / {total}  ·  Module 4",
)

content_slide(
    M4, M4C,
    "Capacity planning rules of thumb",
    "When (not if) to scale each layer",
    "prometheus", "Prometheus + Grafana — your scaling evidence",
    "Restaurant kitchen during a rush",
    "You don't hire more chefs when the dining room is half full. Watch the "
    "metrics, scale the bottleneck, then re-measure. Premature scaling = "
    "expensive empty restaurant.",
    [
        "Kafka: scale partitions when consumer lag stays > 30 s; brokers when "
        "disk fills.",
        "Connect: scale workers when tasks-per-worker > 5 and CPU > 70 %.",
        "MinIO: add drives when used > 70 %; rebalance is online.",
        "Airflow: scale workers when scheduler queue > 0 sustained; switch to "
        "CeleryExecutor before that.",
    ],
    f"Slide 9 / {total}  ·  Module 4",
)

quiz_slide(
    M4, M4C,
    [
        ("Why pre-build day10-app before `compose up -d` in bootstrap?",
         "So compose has the image locally and doesn't waste a network "
         "round-trip + timeout trying to pull it from Docker Hub."),
        ("Where does Alertmanager send firing alerts in this stack?",
         "POST → http://quality-dashboard:5001/webhook/alerts, so DataOps sees "
         "them in the same UI as the metrics."),
        ("What's the single most common cause of 'Debezium silent / 0 events'?",
         "Postgres wal_level is not 'logical' (or the publication/slot wasn't "
         "created). Run 03_enable_wal.sql, restart, re-create the connector."),
        ("What's the cleanest way to handle a poison message in the S3 sink?",
         "Set errors.tolerance=all + errors.deadletterqueue.topic.name — bad "
         "records go to dlq.s3-sink-bronze-* and the sink keeps running."),
        ("How do you make bootstrap.ps1 safe to re-run after a failure?",
         "Idempotent steps: skip downloads if files exist, `compose up -d` is "
         "a no-op for healthy services, schema/connector registration uses PUT."),
    ],
    f"Slide 10 / {total}  ·  Module 4 — Knowledge Check",
)

# =============================================================================
# Closing section  (4 slides)
# =============================================================================
CLOSING_C = NAVY

def closing_what_you_built_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "CLOSING", "What you actually built",
           "Inventory of the running stack — count it, you earned it")

    rows = [
        ("Source feeds",          "3 live + 1 CDC", "OGN · NOAA · EMSC · Postgres"),
        ("Kafka topics",          "8",              "5 streams + 3 CDC + 1 heartbeat"),
        ("Avro schemas registered", "5",            "ogn_position · noaa_observation · noaa_alert · seismic_event · config_change"),
        ("Bronze layout",         "hourly",         "s3://bronze/<topic>/year=…/hour=…/…avro"),
        ("Silver entities",       "4",              "dedupe + late-filter + audit"),
        ("Gold marts",            "4",              "aircraft density · weather snapshot · seismic 24h · region-alert correlation"),
        ("Airflow DAGs",          "5",              "bootstrap + drift + medallion + quality + KPIs"),
        ("Dashboards",            "3",              "quality :5001 · business :5002 · live-map :5003"),
        ("Observability",         "3 services",     "prometheus · alertmanager · grafana"),
        ("Containers",            "22",             "one bridge network · one .env"),
    ]
    n_cols = 3
    tbl = s.shapes.add_table(1 + len(rows), n_cols,
                             Inches(0.3), Inches(1.05),
                             Inches(9.4), Inches(4.4)).table
    for j, h in enumerate(["Component", "Count", "Detail"]):
        cell = tbl.cell(0, j); cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10); r.font.color.rgb = NAVY
                    if j == 1:
                        r.font.bold = True; r.font.color.rgb = ACCENT
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
    footer(s, "Slide 49 / 52  ·  Closing", CLOSING_C)
    return s


def closing_next_steps_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "CLOSING", "Where to take it next",
           "Six concrete extensions — each one builds on what you already have")

    items = [
        ("Add a 4th source feed",
         "ADS-B Exchange or USGS — proves the contract is reusable. "
         "Reuse the ingestor pattern, register a new Avro schema."),
        ("Swap MinIO for S3 (or Azure Blob)",
         "Only the endpoint URL changes. The sink connector, paths, and "
         "DuckDB read URIs stay identical."),
        ("Add CeleryExecutor and a second Airflow worker",
         "When the scheduler queue stays > 0, the LocalExecutor is the bottleneck. "
         "Celery + Redis lets you scale workers horizontally."),
        ("Wire alerts to Slack / PagerDuty",
         "Alertmanager already POSTs to a webhook — add a receiver block for "
         "Slack / PagerDuty / Opsgenie alongside the quality-dashboard one."),
        ("Add a streaming SQL engine (ksqlDB / Flink)",
         "Bronze is already there — point a stream processor at it for "
         "windowed aggregates you can't easily express in the 5-min Hop run."),
        ("Wrap the marts with FastAPI",
         "Same latest.parquet, served as JSON over HTTP — turns dashboards "
         "into headless APIs that any client can consume."),
    ]
    card_w = Inches(4.65); card_h = Inches(1.4); gap = Inches(0.15)
    cols = 2
    start_x = Inches(0.3)
    start_y = Inches(1.05)
    for i, (title, body) in enumerate(items):
        r, c = divmod(i, cols)
        x = start_x + c * (card_w + gap)
        y = start_y + r * (card_h + gap)
        add_rect(s, x, y, Inches(0.18), card_h, ACCENT)
        add_rect(s, x + Inches(0.18), y, card_w - Inches(0.18), card_h,
                 LIGHT, line=ACCENT, line_w=Emu(9525))
        add_text(s, x + Inches(0.3), y + Inches(0.1),
                 card_w - Inches(0.4), Inches(0.35),
                 f"{i+1}.  {title}", size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(0.45),
                 card_w - Inches(0.4), card_h - Inches(0.5),
                 body, size=10, color=GREY)
    footer(s, "Slide 50 / 52  ·  Closing", CLOSING_C)
    return s


def closing_final_quiz_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "CLOSING", "Final knowledge check — mixed",
           "One question from each module + one integration question")
    qa = [
        ("M1 · Why three medallion layers and not just bronze → gold?",
         "Silver is the cheap idempotent join surface — dedupe, late-filter, "
         "audit once so every gold mart can reuse it without re-doing the work."),
        ("M2 · What does the Schema Registry actually enforce on producers?",
         "The compatibility rule for that subject (BACKWARD by default) — a "
         "producer that ships an incompatible schema fails at registration time, "
         "not in production."),
        ("M3 · What's the file the dashboards point at, and why?",
         "latest.parquet — written via atomic rename so the dashboard never "
         "reads a half-finished file, and the URL never changes."),
        ("M4 · After a Kafka restart, what proves the pipeline recovered?",
         "Consumer lag on every S3-sink consumer group returns to ~0, and "
         "bronze hourly files for the current hour resume appearing."),
        ("Integration · A new feed must be onboarded in 1 day — what's the path?",
         "Write an ingestor service (same pattern as ogn/noaa/seismic), "
         "register an Avro schema, add a topic, point the S3-sink at it. "
         "Silver + Gold come along for free if the new entity matches an "
         "existing mart shape."),
    ]
    y = Inches(1.05)
    for i, (q, a) in enumerate(qa, start=1):
        add_text(s, Inches(0.4), y, SW - Inches(0.8), Inches(0.4),
                 f"Q{i}. {q}", size=12, bold=True, color=NAVY)
        add_text(s, Inches(0.6), y + Inches(0.4), SW - Inches(1.0), Inches(0.5),
                 "A. " + a, size=11, color=GREEN)
        y += Inches(0.9)
    footer(s, "Slide 51 / 52  ·  Closing", CLOSING_C)
    return s


def closing_thank_you_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SH, ACCENT)
    add_text(s, Inches(0.5), Inches(0.6), SW - Inches(1.0), Inches(0.5),
             "DAY 10  ·  CAPSTONE", size=14, bold=True,
             color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, Inches(0.5), Inches(1.1), SW - Inches(1.0), Inches(1.0),
             "Thank you — go build it", size=32, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(2.2), SW - Inches(1.0), Inches(0.5),
             "The stack is one  ./bootstrap.ps1  away.",
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
             name="Consolas")

    # URL cheat-sheet
    add_rect(s, Inches(0.5), Inches(3.0), Inches(9.0), Inches(3.4),
             RGBColor(0x07, 0x1C, 0x33), line=ACCENT, line_w=Emu(9525))
    add_text(s, Inches(0.7), Inches(3.1), Inches(8.6), Inches(0.4),
             "URL cheat-sheet (local)", size=13, bold=True,
             color=RGBColor(0xCC, 0xDD, 0xEE))
    urls = [
        ("Airflow",            "http://localhost:8080   (admin / admin)"),
        ("Kafka UI",           "http://localhost:8090"),
        ("Schema Registry",    "http://localhost:8081/subjects"),
        ("Connect REST",       "http://localhost:8083/connectors"),
        ("MinIO console",      "http://localhost:9001   (minioadmin / minioadmin)"),
        ("Quality dashboard",  "http://localhost:5001"),
        ("Business dashboard", "http://localhost:5002"),
        ("Live-map dashboard", "http://localhost:5003"),
        ("Prometheus",         "http://localhost:9090"),
        ("Grafana",            "http://localhost:3000   (admin / admin)"),
    ]
    y = Inches(3.55)
    for label, url in urls:
        add_text(s, Inches(0.8), y, Inches(2.4), Inches(0.25),
                 label, size=10, bold=True, color=WHITE)
        add_text(s, Inches(3.2), y, Inches(6.2), Inches(0.25),
                 url, size=10, color=RGBColor(0xCC, 0xDD, 0xEE),
                 name="Consolas")
        y += Inches(0.27)

    add_text(s, Inches(0.5), SH - Inches(0.55), SW - Inches(1.0), Inches(0.4),
             "Slide 52 / 52  ·  end of deck",
             size=10, color=RGBColor(0xCC, 0xDD, 0xEE),
             align=PP_ALIGN.CENTER)
    return s


closing_what_you_built_slide()
closing_next_steps_slide()
closing_final_quiz_slide()
closing_thank_you_slide()

prs.save(DECK)
print(f"\nsaved: {DECK}")
print(f"total slides now: {len(prs.slides)}")
