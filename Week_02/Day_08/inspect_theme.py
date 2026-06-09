from pptx import Presentation
from lxml import etree
import re

pptx_path = r'C:\Temp\Day_08_Slides.pptx'
prs = Presentation(pptx_path)

master_xml = etree.tostring(prs.slide_master._element, pretty_print=True).decode()

colors = re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', master_xml)
print('Master colors:', sorted(set(colors)))

fonts = re.findall(r'typeface="([^"]+)"', master_xml)
print('Master fonts:', sorted(set(fonts)))

all_colors = set()
for slide in prs.slides:
    xml = etree.tostring(slide._element, pretty_print=True).decode()
    cs = re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml)
    all_colors.update(cs)
print('All slide colors:', sorted(all_colors))

print('\nSlide 1 XML (first 4000):')
print(etree.tostring(prs.slides[0]._element, pretty_print=True).decode()[:4000])

print('\nSlide 2 XML (first 3000):')
print(etree.tostring(prs.slides[1]._element, pretty_print=True).decode()[:3000])

# Print slide backgrounds
print('\n=== SLIDE BACKGROUND INFO ===')
for i, slide in enumerate(prs.slides[:5]):
    bg = slide.background
    fill = bg.fill
    print(f'Slide {i+1} fill type: {fill.type}')
    try:
        xml = etree.tostring(bg._element, pretty_print=True).decode()
        print(f'  BG XML: {xml[:500]}')
    except:
        pass
