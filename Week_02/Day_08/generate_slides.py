"""
Day 8 Airflow Slides Generator
Matches existing Day_08_Slides.pptx design exactly:
  - Background: #F9FBFC
  - Header bar: #0B3D51 (deep teal), height = 457200 EMU (0.5 inch)
  - Accent: #007ACC, #0066CC, #00A78D, #F28C28
  - Major font (titles): Calibri Light
  - Minor font (body): Calibri
  - Slide size: 10 x 5.625 inches (9144000 x 5143500 EMU)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
from copy import deepcopy
from lxml import etree
import textwrap

# ──────────────── COLOUR PALETTE ────────────────
C_BG       = RGBColor(0xF9, 0xFB, 0xFC)   # slide background
C_HEADER   = RGBColor(0x0B, 0x3D, 0x51)   # top bar
C_TITLE    = RGBColor(0x0B, 0x3D, 0x51)   # slide title text
C_ACCENT   = RGBColor(0x00, 0x7A, 0xCC)   # accent blue
C_TEAL     = RGBColor(0x00, 0xA7, 0x8D)   # teal accent
C_ORANGE   = RGBColor(0xF2, 0x8C, 0x28)   # orange highlight
C_BODY     = RGBColor(0x0B, 0x3D, 0x51)   # body text dark
C_BODY2    = RGBColor(0x6A, 0x6A, 0x6A)   # secondary body
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DIVIDER  = RGBColor(0xD0, 0xD7, 0xDE)
C_BOX_BG   = RGBColor(0xE9, 0xF2, 0xF9)   # light blue box
C_BOX_GRN  = RGBColor(0xD0, 0xF0, 0xEB)   # light teal box
C_BOX_ORG  = RGBColor(0xFD, 0xF0, 0xE0)   # light orange box

# ──────────────── SLIDE SIZE ────────────────
W = 9144000   # 10 inches
H = 5143500   # 5.625 inches

HEADER_H = 457200   # 0.5 inch

def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    return prs

def blank_layout(prs):
    # Use layout 6 (blank) if available, else 0
    layouts = prs.slide_layouts
    for l in layouts:
        if l.name.lower() in ('blank', ''):
            return l
    return layouts[6] if len(layouts) > 6 else layouts[0]

# ──────────────── XML HELPERS ────────────────

def rgb_val(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def set_bg(slide, color: RGBColor):
    """Set solid background colour for a slide."""
    # Remove existing bgPr/bgRef
    spTree = slide.shapes._spTree
    cSld = spTree.getparent()
    # Find or create bg element
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    existing = cSld.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if existing is not None:
        cSld.remove(existing)

    bg_xml = f'''<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:bgPr>
        <a:solidFill>
          <a:srgbClr val="{rgb_val(color)}"/>
        </a:solidFill>
        <a:effectLst/>
      </p:bgPr>
    </p:bg>'''
    bg_elem = etree.fromstring(bg_xml)
    # Insert before spTree
    cSld.insert(list(cSld).index(spTree), bg_elem)


def add_rect(slide, x, y, cx, cy, fill_color: RGBColor, line_color=None, line_w=0):
    """Add a filled rectangle shape."""
    from pptx.util import Emu as E
    sp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        E(x), E(y), E(cx), E(cy)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Emu(line_w)
    else:
        sp.line.fill.background()
    return sp


def add_text_box(slide, x, y, cx, cy, text, font_name="Calibri Light",
                 font_size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, italic=False):
    """Add a text box with a single paragraph."""
    from pptx.util import Emu as E
    txBox = slide.shapes.add_textbox(E(x), E(y), E(cx), E(cy))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline_body(slide, x, y, cx, cy, lines, font_name="Calibri",
                       font_size=14, color=None, bullet=True, line_spacing_pt=None):
    """Add a text box with bullet lines."""
    from pptx.util import Emu as E
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(E(x), E(y), E(cx), E(cy))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet:
            # Add bullet using XML
            pPr = p._pPr
            if pPr is None:
                pPr_xml = '<a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" marL="228600" indent="-228600"/>'
                pPr = etree.fromstring(pPr_xml)
                p._p.insert(0, pPr)
            buChar_xml = '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="●"/>'
            pPr.append(etree.fromstring(buChar_xml))
            buClr_xml = f'<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:srgbClr val="{rgb_val(C_ACCENT)}"/></a:buClr>'
            pPr.append(etree.fromstring(buClr_xml))
            buSzPct_xml = '<a:buSzPct xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="80000"/>'
            pPr.append(etree.fromstring(buSzPct_xml))

        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = C_BODY
    return txBox


def header_bar(slide, title=""):
    """Add top dark header bar, optionally with short label text."""
    add_rect(slide, 0, 0, W, HEADER_H, C_HEADER)
    if title:
        add_text_box(slide, 228600, 0, W - 457200, HEADER_H,
                     title, font_name="Calibri Light", font_size=10,
                     color=C_WHITE, align=PP_ALIGN.RIGHT)


def section_pill(slide, label):
    """Small pill label below header indicating section."""
    box_w = 2000000
    box_h = 280000
    # Orange pill rectangle
    add_rect(slide, 228600, HEADER_H + 80000, box_w, box_h, C_ORANGE)
    add_text_box(slide, 228600, HEADER_H + 80000, box_w, box_h,
                 label.upper(), font_name="Calibri", font_size=8,
                 bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_number_footer(slide, num):
    """Slide number in bottom right."""
    add_text_box(slide, W - 500000, H - 280000, 400000, 220000,
                 str(num), font_name="Calibri", font_size=9,
                 color=C_BODY2, align=PP_ALIGN.RIGHT)


def divider_line(slide, y):
    """Horizontal divider line."""
    line_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="99" name="line_divider"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="228600" y="{y}"/><a:ext cx="{W - 457200}" cy="12700"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{rgb_val(C_DIVIDER)}"/></a:solidFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
    </p:sp>'''
    slide.shapes._spTree.append(etree.fromstring(line_xml))


def add_colored_box(slide, x, y, cx, cy, bg_color, text, font_size=11,
                    text_color=None, bold=False, font_name="Calibri"):
    """Rounded-corner box with text."""
    add_rect(slide, x, y, cx, cy, bg_color)
    pad = 80000
    add_text_box(slide, x + pad, y + pad, cx - 2 * pad, cy - 2 * pad,
                 text, font_name=font_name, font_size=font_size,
                 bold=bold, color=text_color or C_BODY)


def add_tag(slide, x, y, text, color=None):
    """Small inline tag/chip."""
    c = color or C_ACCENT
    bg_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="99" name="tag"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="900000" cy="200000"/></a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 16667"/></a:avLst></a:prstGeom>
        <a:solidFill><a:srgbClr val="{c.rgb:06X}"/></a:solidFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr insTex="27000"/><a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r>
          <a:rPr lang="en-IN" sz="700" b="1">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr>
          <a:t>{text}</a:t>
        </a:r></a:p>
      </p:txBody>
    </p:sp>'''
    slide.shapes._spTree.append(etree.fromstring(bg_xml))


# ════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ════════════════════════════════════════════════════════

def build_title_slide(prs, slide_num, title_line1, title_line2, subtitle, section_label="Opening"):
    """Hero title slide (slide 1 style)."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)

    # Top bar
    add_rect(slide, 0, 0, W, HEADER_H, C_HEADER)

    # Left accent strip
    add_rect(slide, 0, HEADER_H, 360000, H - HEADER_H, C_ACCENT)

    # Large title
    title_y = HEADER_H + 400000
    title_h = 1200000
    add_text_box(slide, 500000, title_y, W - 700000, title_h,
                 title_line1, font_name="Calibri Light", font_size=40,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)
    if title_line2:
        add_text_box(slide, 500000, title_y + 1000000, W - 700000, 800000,
                     title_line2, font_name="Calibri Light", font_size=28,
                     bold=False, color=C_ACCENT, align=PP_ALIGN.LEFT)

    # Orange divider
    add_rect(slide, 500000, title_y + 1900000, 1800000, 60000, C_ORANGE)

    # Subtitle
    add_text_box(slide, 500000, title_y + 2100000, W - 700000, 800000,
                 subtitle, font_name="Calibri", font_size=14,
                 color=C_BODY2, align=PP_ALIGN.LEFT)

    # Day tag
    add_colored_box(slide, W - 2400000, HEADER_H + 160000, 2000000, 320000,
                    C_BOX_BG, "Day 8  •  Full-day Session  •  ~8 hours",
                    font_size=10, text_color=C_ACCENT, bold=False)

    slide_number_footer(slide, slide_num)
    return slide


def build_section_divider(prs, slide_num, section_num, section_title, description):
    """Section break slide."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_HEADER)   # Dark background for section breaks

    # White overlay area (right 70%)
    add_rect(slide, int(W * 0.3), 0, int(W * 0.7), H, RGBColor(0xF1, 0xF5, 0xF8))

    # Section number circle area
    add_text_box(slide, 200000, H // 2 - 700000, int(W * 0.28),
                 1400000, str(section_num), font_name="Calibri Light",
                 font_size=80, bold=False, color=C_WHITE,
                 align=PP_ALIGN.CENTER)

    # Section label (right side)
    add_text_box(slide, int(W * 0.32), H // 2 - 600000,
                 int(W * 0.65), 400000,
                 "SECTION", font_name="Calibri", font_size=10,
                 bold=True, color=C_ORANGE, align=PP_ALIGN.LEFT)

    add_text_box(slide, int(W * 0.32), H // 2 - 400000,
                 int(W * 0.65), 900000,
                 section_title, font_name="Calibri Light", font_size=30,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)

    add_rect(slide, int(W * 0.32), H // 2 + 540000, 1200000, 40000, C_ORANGE)

    add_text_box(slide, int(W * 0.32), H // 2 + 640000,
                 int(W * 0.65), 500000,
                 description, font_name="Calibri", font_size=12,
                 color=C_BODY2, align=PP_ALIGN.LEFT)

    slide_number_footer(slide, slide_num)
    return slide


def build_content_slide(prs, slide_num, title, bullets, section_label="",
                        interactive_q=None, interactive_a=None,
                        tip=None, two_col=False, col1_bullets=None, col2_bullets=None,
                        col1_title="", col2_title=""):
    """Standard content slide with title + bullets."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)
    header_bar(slide)

    if section_label:
        section_pill(slide, section_label)

    # Title
    title_y = HEADER_H + 380000
    title_h = 580000
    add_text_box(slide, 228600, title_y, W - 457200, title_h,
                 title, font_name="Calibri Light", font_size=24,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)

    # Orange underline under title
    add_rect(slide, 228600, title_y + title_h, 1200000, 50000, C_ORANGE)

    content_top = title_y + title_h + 100000

    if two_col and col1_bullets is not None:
        # Two-column layout
        col_w = (W - 800000) // 2
        col_h = H - content_top - 350000

        if col1_title:
            add_text_box(slide, 228600, content_top, col_w, 300000,
                         col1_title, font_name="Calibri", font_size=13,
                         bold=True, color=C_ACCENT)
        add_multiline_body(slide, 228600, content_top + 320000, col_w,
                           col_h - 320000, col1_bullets, font_size=11)

        if col2_title:
            add_text_box(slide, 228600 + col_w + 342900, content_top, col_w, 300000,
                         col2_title, font_name="Calibri", font_size=13,
                         bold=True, color=C_TEAL)
        add_multiline_body(slide, 228600 + col_w + 342900, content_top + 320000,
                           col_w, col_h - 320000, col2_bullets, font_size=11)

        divider_line(slide, content_top + col_h // 2)

    else:
        # Single column bullets
        bullet_h = H - content_top - 350000
        if interactive_q or tip:
            bullet_h = H - content_top - 1000000

        add_multiline_body(slide, 228600, content_top + 60000, W - 500000,
                           bullet_h, bullets, font_size=13)

    # Interactive box at bottom
    if interactive_q:
        box_y = H - 900000
        box_h = 600000
        add_rect(slide, 228600, box_y, W - 900000, box_h, C_BOX_BG)
        add_text_box(slide, 350000, box_y + 60000, W - 1100000, 220000,
                     "💡 Quick Check", font_name="Calibri", font_size=9,
                     bold=True, color=C_ACCENT)
        add_text_box(slide, 350000, box_y + 240000, W - 1100000, 200000,
                     f"Q: {interactive_q}", font_name="Calibri", font_size=11,
                     color=C_BODY)
        add_text_box(slide, 350000, box_y + 400000, W - 1100000, 200000,
                     f"A: {interactive_a}", font_name="Calibri", font_size=11,
                     bold=True, color=C_TEAL)

    if tip:
        tip_y = H - 900000
        add_rect(slide, 228600, tip_y, W - 900000, 500000, C_BOX_ORG)
        add_text_box(slide, 350000, tip_y + 60000, W - 1100000, 200000,
                     "Trainer Tip", font_name="Calibri", font_size=9,
                     bold=True, color=C_ORANGE)
        add_text_box(slide, 350000, tip_y + 240000, W - 1100000, 220000,
                     tip, font_name="Calibri", font_size=11, color=C_BODY)

    slide_number_footer(slide, slide_num)
    return slide


def build_diagram_slide(prs, slide_num, title, diagram_items, diagram_type="flow",
                        section_label="", note=""):
    """
    diagram_type: 'flow' (left-to-right boxes), 'arch' (grid boxes), 
                  'table' (two-col table), 'states' (horizontal states)
    """
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)
    header_bar(slide)

    if section_label:
        section_pill(slide, section_label)

    title_y = HEADER_H + 380000
    add_text_box(slide, 228600, title_y, W - 457200, 580000,
                 title, font_name="Calibri Light", font_size=22,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)
    add_rect(slide, 228600, title_y + 580000, 1200000, 50000, C_ORANGE)

    content_top = title_y + 700000
    content_h = H - content_top - 400000

    if diagram_type == "flow":
        n = len(diagram_items)
        box_w = min(1600000, (W - 600000 - (n - 1) * 200000) // n)
        box_h = 600000
        box_y = content_top + (content_h - box_h) // 2
        colors = [C_HEADER, C_ACCENT, C_TEAL, C_ORANGE, C_HEADER, C_ACCENT, C_TEAL]
        for i, item in enumerate(diagram_items):
            bx = 228600 + i * (box_w + 200000)
            add_rect(slide, bx, box_y, box_w, box_h, colors[i % len(colors)])
            add_text_box(slide, bx, box_y, box_w, box_h,
                         item, font_name="Calibri", font_size=11,
                         bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            if i < n - 1:
                arr_x = bx + box_w
                arr_y = box_y + box_h // 2 - 40000
                add_rect(slide, arr_x, arr_y, 200000, 80000, C_ORANGE)

    elif diagram_type == "arch":
        # Grid layout
        cols = 3
        rows = (len(diagram_items) + cols - 1) // cols
        box_w = (W - 600000 - (cols - 1) * 150000) // cols
        box_h = max(400000, (content_h - (rows - 1) * 120000) // rows)
        arch_colors = [C_HEADER, C_ACCENT, C_TEAL, C_ORANGE, C_BOX_BG,
                       C_BOX_GRN, C_BOX_ORG, C_ACCENT]
        txt_colors  = [C_WHITE,  C_WHITE,  C_WHITE,  C_WHITE,  C_BODY,
                       C_BODY,   C_BODY,   C_WHITE]
        for i, (box_title, box_desc) in enumerate(diagram_items):
            r, c = divmod(i, cols)
            bx = 228600 + c * (box_w + 150000)
            by = content_top + r * (box_h + 120000)
            bg = arch_colors[i % len(arch_colors)]
            tc = txt_colors[i % len(txt_colors)]
            add_rect(slide, bx, by, box_w, box_h, bg)
            add_text_box(slide, bx + 80000, by + 60000, box_w - 160000, 250000,
                         box_title, font_name="Calibri", font_size=12,
                         bold=True, color=tc)
            add_text_box(slide, bx + 80000, by + 280000, box_w - 160000, box_h - 340000,
                         box_desc, font_name="Calibri", font_size=9,
                         color=tc if tc == C_WHITE else C_BODY2)

    elif diagram_type == "table":
        # Two-column table
        row_h = 320000
        col_w = (W - 600000) // 2
        header_items, data_rows = diagram_items[0], diagram_items[1:]
        # Header row
        add_rect(slide, 228600, content_top, col_w, row_h, C_HEADER)
        add_text_box(slide, 228600 + 60000, content_top, col_w - 120000, row_h,
                     header_items[0], font_name="Calibri", font_size=12,
                     bold=True, color=C_WHITE)
        add_rect(slide, 228600 + col_w + 60000, content_top, col_w - 60000, row_h, C_ACCENT)
        add_text_box(slide, 228600 + col_w + 120000, content_top, col_w - 180000, row_h,
                     header_items[1], font_name="Calibri", font_size=12,
                     bold=True, color=C_WHITE)
        # Data rows
        for ri, (left, right) in enumerate(data_rows):
            ry = content_top + row_h + ri * row_h
            bg = C_BOX_BG if ri % 2 == 0 else C_BG
            add_rect(slide, 228600, ry, col_w, row_h, bg)
            add_text_box(slide, 228600 + 60000, ry, col_w - 120000, row_h,
                         left, font_name="Calibri", font_size=11,
                         bold=True, color=C_BODY)
            add_rect(slide, 228600 + col_w + 60000, ry, col_w - 60000, row_h,
                     RGBColor(0xF1, 0xF5, 0xF8) if ri % 2 == 0 else C_BG)
            add_text_box(slide, 228600 + col_w + 120000, ry, col_w - 180000, row_h,
                         right, font_name="Calibri", font_size=11, color=C_BODY)

    elif diagram_type == "states":
        # Horizontal state flow (for task states)
        n = len(diagram_items)
        box_w = min(1400000, (W - 600000 - (n - 1) * 180000) // n)
        box_h = 500000
        row1_y = content_top + 100000
        colors = [C_BODY2, C_ACCENT, C_TEAL, RGBColor(0x28, 0xA7, 0x45),
                  RGBColor(0xDC, 0x35, 0x45), C_ORANGE, C_BODY2]
        for i, (state, desc) in enumerate(diagram_items):
            bx = 228600 + i * (box_w + 180000)
            add_rect(slide, bx, row1_y, box_w, box_h, colors[i % len(colors)])
            add_text_box(slide, bx, row1_y, box_w, 280000, state,
                         font_name="Calibri", font_size=11, bold=True,
                         color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text_box(slide, bx + 40000, row1_y + 280000,
                         box_w - 80000, box_h - 280000,
                         desc, font_name="Calibri", font_size=8,
                         color=C_WHITE, align=PP_ALIGN.CENTER)
            if i < n - 1:
                add_rect(slide, bx + box_w, row1_y + (box_h - 60000) // 2,
                         180000, 60000, C_ORANGE)

    if note:
        add_text_box(slide, 228600, H - 350000, W - 500000, 270000,
                     note, font_name="Calibri", font_size=9,
                     color=C_BODY2, italic=True)

    slide_number_footer(slide, slide_num)
    return slide


def build_quiz_slide(prs, slide_num, title, questions_answers, section_label=""):
    """Quiz/exercise slide."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)
    header_bar(slide)

    if section_label:
        section_pill(slide, section_label)

    title_y = HEADER_H + 380000
    add_text_box(slide, 228600, title_y, W - 457200, 580000,
                 title, font_name="Calibri Light", font_size=22,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)
    add_rect(slide, 228600, title_y + 580000, 1200000, 50000, C_ORANGE)

    content_top = title_y + 700000
    q_spacing = (H - content_top - 350000) // max(len(questions_answers), 1)
    q_h = min(q_spacing - 50000, 750000)

    for i, (q, opts, ans) in enumerate(questions_answers):
        qy = content_top + i * (q_spacing)
        add_rect(slide, 228600, qy, W - 500000, q_h,
                 C_BOX_BG if i % 2 == 0 else C_BOX_GRN)
        add_text_box(slide, 350000, qy + 60000, W - 750000, 220000,
                     f"{i+1}. {q}", font_name="Calibri", font_size=12,
                     bold=True, color=C_TITLE)
        add_text_box(slide, 350000, qy + 260000, W - 750000, 220000,
                     opts, font_name="Calibri", font_size=11, color=C_BODY)
        add_text_box(slide, 350000, qy + 460000, W - 750000, 200000,
                     f"✓ {ans}", font_name="Calibri", font_size=11,
                     bold=True, color=C_TEAL)

    slide_number_footer(slide, slide_num)
    return slide


def build_lab_slide(prs, slide_num, title, steps, code_snippet="",
                    section_label="Lab", objective=""):
    """Hands-on lab slide."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)
    header_bar(slide)
    section_pill(slide, section_label)

    title_y = HEADER_H + 380000
    add_text_box(slide, 228600, title_y, W - 457200, 500000,
                 title, font_name="Calibri Light", font_size=22,
                 bold=False, color=C_TITLE, align=PP_ALIGN.LEFT)
    add_rect(slide, 228600, title_y + 500000, 1200000, 50000, C_ORANGE)

    if objective:
        add_rect(slide, 228600, title_y + 600000, W - 500000, 280000, C_TEAL)
        add_text_box(slide, 350000, title_y + 600000, W - 700000, 280000,
                     f"Objective: {objective}", font_name="Calibri", font_size=11,
                     bold=False, color=C_WHITE)

    content_top = title_y + (940000 if objective else 650000)

    if code_snippet:
        # Split: steps left, code right
        col_w = int((W - 600000) * 0.42)
        code_w = int((W - 600000) * 0.55)
        col_h = H - content_top - 350000

        add_multiline_body(slide, 228600, content_top, col_w, col_h,
                           steps, font_size=11)

        # Code box - dark background
        code_x = 228600 + col_w + 200000
        add_rect(slide, code_x, content_top, code_w, col_h,
                 RGBColor(0x1E, 0x2A, 0x38))
        # Code label
        add_text_box(slide, code_x + 80000, content_top + 60000,
                     code_w - 160000, 200000,
                     "CODE", font_name="Calibri", font_size=8,
                     bold=True, color=C_TEAL)
        add_text_box(slide, code_x + 80000, content_top + 240000,
                     code_w - 160000, col_h - 300000,
                     code_snippet, font_name="Courier New", font_size=9,
                     color=RGBColor(0xA8, 0xD8, 0xAA))
    else:
        col_h = H - content_top - 350000
        add_multiline_body(slide, 228600, content_top, W - 500000, col_h,
                           steps, font_size=12)

    slide_number_footer(slide, slide_num)
    return slide


def build_recap_slide(prs, slide_num, title, recap_items):
    """End-of-day recap/summary slide."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_BG)

    # Full-width dark top
    add_rect(slide, 0, 0, W, HEADER_H * 2, C_HEADER)
    add_text_box(slide, 228600, 0, W - 457200, HEADER_H * 2,
                 title, font_name="Calibri Light", font_size=22,
                 bold=False, color=C_WHITE, align=PP_ALIGN.LEFT)

    content_top = HEADER_H * 2 + 200000
    cols = 2
    rows = (len(recap_items) + cols - 1) // cols
    box_w = (W - 600000 - (cols - 1) * 200000) // cols
    box_h = max(400000, (H - content_top - 300000 - (rows - 1) * 150000) // rows)
    colors = [C_ACCENT, C_TEAL, C_ORANGE, C_HEADER, C_ACCENT, C_TEAL]

    for i, (icon_label, text) in enumerate(recap_items):
        r, c = divmod(i, cols)
        bx = 228600 + c * (box_w + 200000)
        by = content_top + r * (box_h + 150000)
        add_rect(slide, bx, by, box_w, box_h, colors[i % len(colors)])
        add_text_box(slide, bx + 80000, by + 60000, box_w - 160000, 260000,
                     icon_label, font_name="Calibri", font_size=12,
                     bold=True, color=C_WHITE)
        add_text_box(slide, bx + 80000, by + 300000, box_w - 160000,
                     box_h - 360000,
                     text, font_name="Calibri", font_size=10, color=C_WHITE)

    slide_number_footer(slide, slide_num)
    return slide


def build_break_slide(prs, slide_num, break_type, duration, time_hint):
    """Break / lunch slide."""
    layout = blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_HEADER)

    add_rect(slide, 0, 0, W, HEADER_H, C_ORANGE)

    add_text_box(slide, 0, H // 2 - 700000, W, 800000,
                 break_type, font_name="Calibri Light", font_size=52,
                 bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0, H // 2 + 200000, W, 500000,
                 duration, font_name="Calibri Light", font_size=28,
                 bold=False, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0, H // 2 + 700000, W, 400000,
                 time_hint, font_name="Calibri", font_size=14,
                 color=RGBColor(0xA8, 0xD8, 0xAA), align=PP_ALIGN.CENTER)

    slide_number_footer(slide, slide_num)
    return slide


# ════════════════════════════════════════════════════════
#  ALL SLIDE DATA
# ════════════════════════════════════════════════════════

def generate_all_slides():
    prs = new_prs()

    # ── S1: Title ──────────────────────────────────────────
    build_title_slide(prs, 1,
        "Day 8",
        "Workflow Orchestration with Apache Airflow",
        "Design · Schedule · Monitor · Integrate with Kafka & Hop",
        "Opening")

    # ── S2: Learning Outcomes ─────────────────────────────
    build_content_slide(prs, 2, "Learning Outcomes for Today",
        ["Explain Airflow architecture and each component's role",
         "Read and write a DAG with tasks and dependencies",
         "Choose the right operator for each type of work",
         "Schedule, retry, and monitor workflows",
         "Connect Airflow with Kafka topics and Apache Hop pipelines"],
        section_label="Opening",
        interactive_q="By end of day, what tool will you use to schedule a Hop pipeline?",
        interactive_a="Apache Airflow")

    # ── S3: Business Problem ──────────────────────────────
    build_content_slide(prs, 3, "The Business Problem Before Airflow",
        ["A retail company receives orders every night",
         "Script 1 downloads files → Script 2 cleans data → Script 3 loads warehouse → Script 4 sends report",
         "Transform script fails at 2 AM — nobody knows until the CEO report is empty",
         "Missing: scheduling, dependency management, retries, logging, monitoring",
         "One script failure silently ruins the whole pipeline"],
        section_label="Opening",
        interactive_q="What is worse: script failure or not knowing it failed?",
        interactive_a="Not knowing — the business may use wrong or missing data")

    # ── SECTION 1: Introduction ───────────────────────────
    build_section_divider(prs, 4, 1, "Introduction", "From manual scripts to orchestrated workflows")

    # ── S4: Workflow Orchestration ────────────────────────
    build_content_slide(prs, 5, "What is Workflow Orchestration?",
        ["Arranging many steps in the right order, at the right time",
         "Monitoring each step and recovering from failures",
         "Making dependencies explicit — step 2 only runs after step 1 succeeds",
         "Think of a wedding planner: does not cook food, but coordinates everything",
         "Airflow is the orchestrator — not the worker"],
        section_label="Introduction",
        interactive_q="Fill in the blank: Airflow defines workflows as Python ______.",
        interactive_a="code")

    # ── S5: Why Cron Is Not Enough ────────────────────────
    build_content_slide(prs, 6, "Why Cron Is Not Enough",
        ["Cron can start a command at a time — that is all",
         "No native dependency management between jobs",
         "cron starts load.sh at 09:00 even if extract.sh failed at 08:55",
         "No visual monitoring, no built-in retries, no shared logs",
         "When scripts grow from 5 to 50, cron becomes unmanageable"],
        section_label="Introduction",
        interactive_q="Cron is good for simple single jobs. True or false for complex pipelines?",
        interactive_a="False — cron has no dependency awareness or retry logic")

    # ── S6: Why Shell Scripts Become Risky ───────────────
    build_content_slide(prs, 7, "Why Shell Scripts Become Risky",
        ["A script can do work. A workflow platform explains and controls the work",
         "Large script chains are difficult to test, monitor, and hand over",
         "No standard way to retry only the failed step",
         "Credentials often hardcoded — security risk",
         "No single view: which scripts ran today? Which failed?"],
        section_label="Introduction",
        tip="Ask: 'How would a new engineer know which script to restart after a failure?'")

    # ── S7: Airflow in One Sentence ───────────────────────
    build_content_slide(prs, 8, "Airflow in One Sentence",
        ["Apache Airflow lets teams define, schedule, and monitor workflows as Python code",
         "Define = write the workflow structure once",
         "Schedule = decide when it runs (daily, hourly, on event)",
         "Monitor = see status, logs, retries from a single UI",
         "Python code = version controlled, reviewable, testable"],
        section_label="Introduction",
        interactive_q="Airflow workflows are defined as ______.",
        interactive_a="Python code")

    # ── S8: Not a Data Processing Engine ─────────────────
    build_content_slide(prs, 9, "Airflow is NOT a Data Processing Engine",
        ["Airflow coordinates work — it does not process big data directly",
         "Airflow is NOT Spark, Kafka, Hop, Snowflake, or a database",
         "Airflow CAN call a Spark job, Hop pipeline, SQL query, or Kafka check",
         "Heavy data processing should happen inside the called tools",
         "Running a pandas loop on 10GB inside a PythonOperator is bad practice"],
        section_label="Introduction",
        interactive_q="Should Airflow process 10GB of raw CSV data directly?",
        interactive_a="No — use Hop, Spark, or Snowflake for heavy processing. Airflow orchestrates.")

    # ── S9: Where Airflow Fits ────────────────────────────
    build_content_slide(prs, 10, "Where Airflow Fits in Data Engineering",
        ["Airflow sits above data systems — like a conductor above an orchestra",
         "Controls: ETL, ELT, data quality checks, notifications, housekeeping",
         "Calls: databases, warehouses, file systems, APIs, containers",
         "Does NOT replace: Kafka (streaming), Hop (visual ETL), Spark (compute)",
         "Each tool has one job — Airflow is the coordinator"],
        section_label="Introduction")

    # ── S10: Who Uses Airflow ─────────────────────────────
    build_content_slide(prs, 11, "Who Uses Airflow?",
        ["Data engineers: daily warehouse refresh, pipeline scheduling",
         "Analytics engineers: dbt model runs, data quality checks",
         "ML engineers: model retraining pipelines, feature engineering",
         "DevOps/Platform teams: infrastructure automation jobs",
         "Reporting teams: scheduled report generation and distribution"],
        section_label="Introduction",
        interactive_q="Which team would use Airflow to automate nightly ML model retraining?",
        interactive_a="ML engineers — they use Airflow to schedule and monitor training pipelines")

    # ── S11: Airflow Benefits ─────────────────────────────
    build_content_slide(prs, 12, "Airflow Benefits Overview",
        ["Code-based workflows: version controlled, tested, reviewed like software",
         "Dependency management: task B only runs if task A succeeded",
         "Scheduling: cron, data-driven events, manual trigger",
         "Retries & timeouts: automatic recovery from transient failures",
         "Web UI: visual graph, logs, task history, trigger buttons",
         "Extensibility: 1000+ providers for databases, cloud, APIs"],
        section_label="Introduction",
        tip="Connect each benefit to a pain from Slide 3 (the retail pipeline story)")

    # ── SECTION 2: Architecture ───────────────────────────
    build_section_divider(prs, 13, 2, "Architecture", "The components that make Airflow work together")

    # ── S12: Architecture Mental Model ────────────────────
    build_diagram_slide(prs, 14, "Airflow Architecture: The Control Room",
        [("DAG Files", "Python workflow definitions"),
         ("DAG Processor", "Parses .py files, serialises to DB"),
         ("Metadata DB", "Stores all states, history, variables"),
         ("Scheduler", "Decides WHEN to run each task"),
         ("Executor", "Decides HOW to run: local / Celery / K8s"),
         ("Worker", "Actually executes the task code"),
         ("Web UI", "Visual control room — graphs, logs, triggers"),
         ("Triggerer", "Async waiting for deferrable sensors"),
         ("Logs", "Per-task execution output, accessible from UI")],
        diagram_type="arch",
        section_label="Architecture",
        note="Ref: airflow.apache.org/docs/apache-airflow/stable/core-concepts/overview.html")

    # ── S13: DAG Folder ───────────────────────────────────
    build_content_slide(prs, 15, "The DAG Folder",
        ["The dags/ folder contains Python files describing workflows",
         "Airflow reads them to discover what workflows exist",
         "Placing a file here is like submitting a route plan — it does not run immediately",
         "Files are re-parsed regularly — DAG Processor notices changes",
         "Syntax errors in any file cause import errors visible in the UI"],
        section_label="Architecture",
        interactive_q="If you place a Python file in the dags/ folder, does the DAG run immediately?",
        interactive_a="No — it is discovered and scheduled; it only runs when the schedule triggers")

    # ── S14: DAG Processor ────────────────────────────────
    build_content_slide(prs, 16, "The DAG Processor",
        ["Reads Python DAG files from the dags/ folder",
         "Parses them and stores serialised DAG structure in the Metadata DB",
         "The Scheduler reads from the DB — not directly from files",
         "Multiple processors can run in parallel for large numbers of DAGs",
         "Parse errors appear as import errors in the Web UI"],
        section_label="Architecture",
        tip="Say: 'Before the scheduler can schedule, someone must read the Python files and understand their structure'")

    # ── S15: Scheduler ────────────────────────────────────
    build_content_slide(prs, 17, "The Scheduler",
        ["Brain of Airflow — decides when DAG runs and task instances are created",
         "Asks: Is it time? Are upstream tasks done? Are there free slots?",
         "Reads from the Metadata DB, creates task instances, hands off to Executor",
         "Runs continuously — checks for newly eligible tasks every few seconds",
         "Station master analogy: decides when trains depart"],
        section_label="Architecture",
        interactive_q="Scheduler = ?  A) Stores states  B) Decides when tasks run  C) ETL tool",
        interactive_a="B) Decides when tasks run")

    # ── S16: Executor ─────────────────────────────────────
    build_content_slide(prs, 18, "The Executor",
        ["The Executor is the strategy for how tasks are physically run",
         "LocalExecutor: runs tasks as subprocesses on the scheduler machine",
         "CeleryExecutor: queues tasks to distributed workers via broker (Redis/RabbitMQ)",
         "KubernetesExecutor: spins up a fresh pod per task",
         "SequentialExecutor: one task at a time — only for testing"],
        section_label="Architecture",
        interactive_q="CeleryExecutor is best for: A) Learning  B) Distributed production scale",
        interactive_a="B — Celery lets tasks fan out to many worker machines")

    # ── S17: Workers ──────────────────────────────────────
    build_content_slide(prs, 19, "Workers",
        ["Workers execute the actual task code in distributed setups",
         "Delivery driver analogy: Scheduler decides the package goes; Worker delivers",
         "With LocalExecutor, the scheduler machine IS the worker",
         "Workers need access to DAG files and required Python packages",
         "Multiple workers allow parallel task execution at scale"],
        section_label="Architecture",
        tip="For the Docker Compose lab we use LocalExecutor — one machine does everything")

    # ── S18: Metadata Database ────────────────────────────
    build_content_slide(prs, 20, "The Metadata Database",
        ["PostgreSQL (or MySQL) database that stores ALL Airflow state",
         "Stores: DAG run history, task instance states, variables, connections",
         "Without it, Airflow forgets everything on restart — it is the memory",
         "Ledger analogy: every run outcome written here permanently",
         "In the lab: PostgreSQL container managed by Docker Compose"],
        section_label="Architecture",
        interactive_q="Which component stores task run history and task states?",
        interactive_a="Metadata Database (usually PostgreSQL)")

    # ── S19: Webserver and UI ─────────────────────────────
    build_content_slide(prs, 21, "Webserver and the UI",
        ["Provides the browser-based control room — http://localhost:8080 in the lab",
         "Shows: DAG list, graph view, task states, logs, calendar, statistics",
         "Buttons: trigger run, pause DAG, clear task, mark success/failed",
         "The UI is not just pretty — it is the primary debugging window",
         "Reads from Metadata DB — it does not run tasks"],
        section_label="Architecture",
        interactive_q="The Airflow UI typically runs on port ______ in local setups.",
        interactive_a="8080")

    # ── S20: Logs ─────────────────────────────────────────
    build_content_slide(prs, 22, "Task Logs",
        ["Every task instance writes detailed execution logs",
         "Logs show: task start time, print statements, errors, exit code",
         "Accessible from the UI: click a task → Logs tab",
         "Every production support incident starts with logs",
         "Can be stored locally or in remote storage (S3, GCS, Azure Blob)"],
        section_label="Architecture",
        tip="Teach students: DAG list → click DAG → Graph → click task → Logs")

    # ── S21: Triggerer ────────────────────────────────────
    build_content_slide(prs, 23, "The Triggerer",
        ["Handles deferrable tasks that wait asynchronously",
         "Instead of a worker blocking for hours, a lightweight triggerer watches",
         "Worker releases its slot when task defers; triggerer monitors the condition",
         "When condition met, triggerer re-queues the task on a worker",
         "Watchman analogy: does not occupy a seat, just watches the gate"],
        section_label="Architecture")

    # ── S22: Message Broker ───────────────────────────────
    build_content_slide(prs, 24, "Message Broker (Celery Setup)",
        ["In distributed setups, a broker queues tasks between Scheduler and Workers",
         "Common choices: Redis (simple, fast) or RabbitMQ (feature-rich)",
         "Food delivery analogy: orders sit in queue between restaurant and drivers",
         "Scheduler puts tasks on the queue; Workers pick them up when ready",
         "In the Docker Compose lab: Redis is the broker"],
        section_label="Architecture",
        interactive_q="Which component sits between Scheduler and Workers in Celery setup?",
        interactive_a="Message Broker (Redis or RabbitMQ)")

    # ── S23: Minimum Installation ─────────────────────────
    build_content_slide(prs, 25, "Minimum Airflow Installation",
        ["Required for any Airflow setup: Scheduler, DAG Processor, Web UI, Metadata DB",
         "Optional (for scale): Workers, Broker, Triggerer, multiple Scheduler replicas",
         "For local learning: SequentialExecutor or LocalExecutor is sufficient",
         "The official Docker Compose quick-start bundles all services for evaluation",
         "Never present the quick-start compose as production architecture"],
        section_label="Architecture",
        tip="Walk through the architecture diagram slowly — do not mention every advanced service until students understand the minimum")

    # ── S24: Single Machine ───────────────────────────────
    build_content_slide(prs, 26, "Single-Machine Architecture",
        ["All services live in one container / machine",
         "Scheduler, UI, DB, and executor all colocated",
         "Good for: learning, development, testing small workflows",
         "Not suitable for: large-scale production with 1000s of tasks",
         "All services run in one house — if house is crowded, everything slows"],
        section_label="Architecture")

    # ── S25: Celery Architecture ──────────────────────────
    build_diagram_slide(prs, 27, "Distributed Celery Architecture",
        ["Scheduler", "Broker\n(Redis/RabbitMQ)", "Worker 1", "Worker 2", "Worker 3", "Metadata DB"],
        diagram_type="flow",
        section_label="Architecture",
        note="Tasks fan out from Scheduler → Broker → multiple Workers in parallel")

    # ── S26: Kubernetes Architecture ──────────────────────
    build_content_slide(prs, 28, "Kubernetes Architecture",
        ["KubernetesExecutor runs each task as a separate ephemeral pod",
         "Pod = temporary small container room for one task only",
         "After task finishes, pod is deleted — elastic and isolated",
         "Great for heterogeneous workloads with different dependencies",
         "More complex setup — not needed for Day 8 lab"],
        section_label="Architecture",
        interactive_q="What is a Kubernetes pod in the context of Airflow task execution?",
        interactive_a="A temporary container that runs one task and is deleted when done")

    # ── S27: Execution Flow ───────────────────────────────
    build_diagram_slide(prs, 29, "Full Execution Flow: DAG File → UI",
        ["1. DAG file\nplaced in\ndags/", "2. DAG\nProcessor\nparses it", "3. Metadata\nDB stores\nstructure",
         "4. Scheduler\ncreates task\ninstances", "5. Executor\ndispatches\ntask", "6. Worker\nexecutes\ncode",
         "7. Logs\nwritten", "8. State\nupdated\nin DB", "9. Web UI\nshows result"],
        diagram_type="flow",
        section_label="Architecture",
        note="Make students repeat this flow aloud — it is the core architecture memory exercise")

    # ── S28: Architecture Recap ───────────────────────────
    build_diagram_slide(prs, 30, "Architecture Recap: Control Room Roles",
        [("Scheduler", "Decides WHEN\n= Station master"),
         ("DAG Processor", "Reads files\n= File reader"),
         ("Metadata DB", "Remembers all\n= Ledger/register"),
         ("Executor", "Dispatch strategy\n= Transport method"),
         ("Worker", "Does the work\n= Delivery driver"),
         ("Web UI", "Shows status\n= Control room screen"),
         ("Triggerer", "Async watchman\n= Gate monitor"),
         ("Broker", "Task queue\n= Order queue"),
         ("Logs", "Evidence trail\n= Black box recorder")],
        diagram_type="arch",
        section_label="Architecture")

    # ── SECTION 3: DAGs and Tasks ─────────────────────────
    build_section_divider(prs, 31, 3, "DAGs and Tasks", "Directed Acyclic Graphs — the recipe for your workflow")

    # ── S29: What is a DAG ────────────────────────────────
    build_content_slide(prs, 32, "What is a DAG?",
        ["DAG = Directed Acyclic Graph — a workflow made of tasks and dependencies",
         "Directed = arrows have direction (A → B means B comes after A)",
         "Acyclic = no circular loops (A cannot depend on C if C depends on A)",
         "Graph = tasks are nodes, dependencies are edges",
         "Draw it: A → B → C (extract → transform → load)"],
        section_label="DAGs & Tasks",
        interactive_q="DAG stands for Directed ______ Graph.",
        interactive_a="Acyclic")

    # ── S30: Directed ─────────────────────────────────────
    build_content_slide(prs, 33, "Directed Means Order Matters",
        ["Arrows have direction — you cannot reverse them without rewriting the DAG",
         "Cooking analogy: wash vegetables → cut → cook (cannot cook before cutting)",
         "In data: extract raw data → validate → transform → load (strict order)",
         "If task B needs task A's output, B must come after A",
         "The >> operator in Python expresses this: task_a >> task_b"],
        section_label="DAGs & Tasks")

    # ── S31: Acyclic ──────────────────────────────────────
    build_content_slide(prs, 34, "Acyclic Means No Circular Dependency",
        ["A cycle would create an impossible dependency loop",
         "A → B → C → A: who starts first? Nobody can — it is a deadlock",
         "Airflow will reject a DAG with a cycle at parse time",
         "This is why DAGs are safe to schedule without infinite loops",
         "Loops inside task code are fine — loops in the DAG structure are not"],
        section_label="DAGs & Tasks",
        interactive_q="If A depends on C, and C depends on A, what happens?",
        interactive_a="Deadlock — nobody can start. Airflow rejects this as a cycle error.")

    # ── S32: Graph ────────────────────────────────────────
    build_diagram_slide(prs, 35, "Graph Means a Visual Workflow",
        ["Extract\nraw data", "Validate\nrows", "Transform\ncleaned", "Load\nwarehouse", "Send\nreport"],
        diagram_type="flow",
        section_label="DAGs & Tasks",
        note="The Python code becomes this visual map in the Airflow UI Graph view")

    # ── S33: DAG File is Python Config ────────────────────
    build_content_slide(prs, 36, "DAG File: Python Configuration, Not Processing",
        ["DAG file is parsed frequently — keep it lightweight",
         "At import time: describe the recipe (define tasks, dependencies, schedule)",
         "At task runtime: cook the food (do actual data work inside tasks)",
         "Never read files, call APIs, or query DBs at the top level of a DAG file",
         "Heavy top-level code slows the DAG Processor and can cause timeouts"],
        section_label="DAGs & Tasks",
        tip="Rule of thumb: if it runs when you import the file, it is too heavy for the DAG file")

    # ── S34: DAG ID ───────────────────────────────────────
    build_content_slide(prs, 37, "DAG ID and Naming",
        ["dag_id uniquely identifies the workflow in Airflow — like a filename",
         "Good names: daily_customer_etl, medallion_bronze_load, kafka_orders_hourly",
         "Bad names: dag1, test_final_latest, my_dag_v3_copy",
         "The dag_id appears in the UI, logs, CLI, and API — make it descriptive",
         "Changing the dag_id loses historical run data — choose names carefully"],
        section_label="DAGs & Tasks",
        interactive_q="Which dag_id is better: 'dag1' or 'daily_sales_etl'?",
        interactive_a="daily_sales_etl — descriptive, readable, searchable in logs and UI")

    # ── S35: Start Date and Schedule ─────────────────────
    build_content_slide(prs, 38, "Start Date and Schedule",
        ["start_date: the earliest date for which a DAG run is eligible",
         "schedule: how often the DAG runs (@daily, @hourly, cron expression, or None)",
         "start_date is NOT 'run now' — it participates in scheduling logic",
         "First scheduled run happens at start_date + one schedule interval",
         "Use a fixed past date like datetime(2024,1,1) for new DAGs — not relative dates"],
        section_label="DAGs & Tasks",
        tip="Common mistake: using datetime.now() as start_date — this causes unpredictable scheduling")

    # ── S36: Catchup ──────────────────────────────────────
    build_content_slide(prs, 39, "Catchup — Be Careful",
        ["catchup=True: Airflow creates missed DAG runs between start_date and today",
         "If start_date is Jan 1 and today is Apr 1, a daily DAG creates ~90 runs",
         "This can flood the executor and overwhelm external systems",
         "catchup=False (recommended default): only schedule from today forward",
         "For deliberate historical reloads, use controlled backfill with limits"],
        section_label="DAGs & Tasks",
        interactive_q="You set start_date=Jan 1 and catchup=True. Today is April 1. What happens?",
        interactive_a="~90 daily DAG runs are created immediately — this can overload the system")

    # ── S37: Default Args ─────────────────────────────────
    build_content_slide(prs, 40, "default_args — Shared Task Settings",
        ["default_args is a dict of settings applied to all tasks in the DAG",
         "Common args: owner, retries, retry_delay, email_on_failure, depends_on_past",
         "School rules analogy: all students follow the default; exceptions override it",
         "A task can override default_args with its own parameter values",
         "Centralises repetitive settings — change once, applies everywhere"],
        section_label="DAGs & Tasks")

    # ── S38: Task ─────────────────────────────────────────
    build_content_slide(prs, 41, "What is a Task?",
        ["A task is one unit of work inside a DAG",
         "Each task has a unique task_id within its DAG",
         "Examples: extract_orders, validate_rows, load_warehouse, send_email",
         "A task is created by instantiating an Operator inside a DAG context",
         "Tasks are connected with >> or << dependency operators"],
        section_label="DAGs & Tasks",
        interactive_q="Name 3 tasks in an ETL pipeline.",
        interactive_a="extract, transform (or validate), load — or any equivalent meaningful names")

    # ── S39: Task Instance ────────────────────────────────
    build_content_slide(prs, 42, "Task Instance vs Task",
        ["Task: the definition (same thing every day — e.g., load_warehouse)",
         "Task Instance: one specific execution — load_warehouse for 2024-01-15",
         "Every scheduled run creates a new set of task instances",
         "Think: class (Task) vs object (Task Instance) — same code, different date",
         "Task instances have their own state, logs, and execution time"],
        section_label="DAGs & Tasks",
        tip="Use analogy: a school timetable is the task; each actual class period is a task instance")

    # ── S40: Task States ──────────────────────────────────
    build_diagram_slide(prs, 43, "Task States",
        [("queued", "Ready to run,\nawaiting worker"),
         ("running", "Currently\nexecuting"),
         ("success", "Finished\nwithout error"),
         ("failed", "Raised an\nexception"),
         ("retry", "Will retry\nafter delay"),
         ("skipped", "Branching chose\nanother path"),
         ("upstream_failed", "Upstream task\nfailed")],
        diagram_type="states",
        section_label="DAGs & Tasks",
        note="Colours in UI: green=success, red=failed, yellow=running, orange=retry, grey=skipped")

    # ── S41: Dependencies ─────────────────────────────────
    build_content_slide(prs, 44, "Expressing Dependencies with >> and <<",
        ["task_a >> task_b: task_b waits for task_a to succeed",
         "task_a << task_b: task_a waits for task_b (less common, read right-to-left)",
         "Chain: extract >> transform >> load",
         "Fan-out: extract >> [validate_customers, validate_orders]",
         "[validate_customers, validate_orders] >> load  (fan-in)"],
        section_label="DAGs & Tasks",
        interactive_q="Write the dependency: extract must finish before transform, transform before load.",
        interactive_a="extract >> transform >> load")

    # ── S42: Linear ETL DAG ───────────────────────────────
    build_diagram_slide(prs, 45, "Pattern: Linear ETL DAG",
        ["extract\nsource data", "validate\nrow counts", "transform\ncleaned data",
         "load\nwarehouse", "send\nreport"],
        diagram_type="flow",
        section_label="DAGs & Tasks",
        note="Most real workflows start with this linear shape. Master it before learning fan-out.")

    # ── S43: Fan-out ──────────────────────────────────────
    build_content_slide(prs, 46, "Pattern: Fan-Out — Parallel Tasks",
        ["One upstream task triggers multiple independent downstream tasks",
         "Example: extract_all >> [validate_customers, validate_orders, validate_payments]",
         "All three validation tasks run in parallel — faster than sequential",
         "Each branch is independent — a failure in one does not automatically stop others (depends on trigger rule)",
         "Use when downstream tasks share no data with each other"],
        section_label="DAGs & Tasks",
        interactive_q="Why is fan-out useful?",
        interactive_a="Independent tasks run in parallel, reducing total pipeline time")

    # ── S44: Fan-in ───────────────────────────────────────
    build_content_slide(prs, 47, "Pattern: Fan-In — Join Branches",
        ["Multiple branches converge into one downstream task",
         "Example: [validate_customers, validate_orders] >> load_warehouse",
         "load_warehouse only runs when ALL validations succeed (with all_success trigger rule)",
         "Useful for aggregating or combining results from parallel branches",
         "Fan-out + Fan-in is a diamond pattern — common in real pipelines"],
        section_label="DAGs & Tasks")

    # ── S45: Branching ────────────────────────────────────
    build_content_slide(prs, 48, "Pattern: Branching",
        ["BranchPythonOperator returns the task_id of the branch to execute",
         "Other branches are automatically skipped",
         "Example: if file has rows → load it; if empty → skip load, send alert",
         "Skipped tasks appear grey in the UI — not failures",
         "Downstream tasks of skipped branches may need trigger_rule=none_failed"],
        section_label="DAGs & Tasks",
        interactive_q="If a BranchPythonOperator returns 'load_data', what happens to 'send_alert'?",
        interactive_a="send_alert is skipped (shown grey in UI)")

    # ── S46: DAG Recap ────────────────────────────────────
    build_content_slide(prs, 49, "DAG Recap: Recipe Analogy",
        ["DAG = recipe (defines the complete workflow)",
         "Task = one step in the recipe (e.g., 'chop vegetables')",
         "Dependency = cooking order (cannot cook before chopping)",
         "DAG Run = one execution of the recipe (making the dish on Monday)",
         "Task Instance = one step executed on a specific run"],
        section_label="DAGs & Tasks",
        interactive_q="Explain a DAG using any food/cooking analogy.",
        interactive_a="The recipe is the DAG. Each cooking step is a task. You cannot fry before washing.")

    # ── SECTION 4: Operators, Sensors, Hooks ─────────────
    build_section_divider(prs, 50, 4, "Operators, Sensors & Hooks",
        "The building blocks of tasks — choosing the right tool for each job")

    # ── S47: What is an Operator ──────────────────────────
    build_content_slide(prs, 51, "What is an Operator?",
        ["An operator is a reusable template for a kind of task",
         "Task is the job in this DAG. Operator is the job type.",
         "BashOperator: runs shell commands",
         "PythonOperator: runs a Python function",
         "1000+ providers offer operators for databases, cloud, APIs, Spark, etc."],
        section_label="Operators & Sensors",
        interactive_q="Which operator would you use to run a shell script?",
        interactive_a="BashOperator")

    # ── S48: Operator vs Task ─────────────────────────────
    build_content_slide(prs, 52, "Operator vs Task",
        ["Operator: the class/template (job description)",
         "Task: the instance created inside a DAG (the actual assigned job)",
         "BashOperator(task_id='extract', bash_command='...') creates a task called extract",
         "Multiple tasks can use the same operator type with different parameters",
         "Same class, different instances — like stamping different orders from one template"],
        section_label="Operators & Sensors",
        tip="Ask: 'Is PythonOperator a task or an operator?' — Answer: both, depending on context")

    # ── S49–S55: Operator Types ────────────────────────────
    build_diagram_slide(prs, 53, "Common Operators Reference",
        [("BashOperator", "Runs shell commands\nor scripts\n→ hop-run, scripts"),
         ("PythonOperator", "Runs a Python\ncallable/function\n→ validation, API calls"),
         ("SQLExecuteQuery\nOperator", "Executes SQL\nvia a connection\n→ warehouse transforms"),
         ("HttpOperator", "Calls REST APIs\n→ trigger/check\nexternal services"),
         ("EmailOperator", "Sends email\n→ success/failure\nnotifications"),
         ("DockerOperator", "Runs Docker\ncontainer\n→ isolated workloads"),
         ("FileSensor", "Waits for file\nto exist\n→ landing file check"),
         ("ExternalTask\nSensor", "Waits for task\nin another DAG\n→ cross-DAG deps"),
         ("BranchPython\nOperator", "Conditional\nbranching\n→ if/else in pipeline")],
        diagram_type="arch",
        section_label="Operators & Sensors",
        note="Ref: airflow.apache.org/docs/apache-airflow/stable/core-concepts/tasks.html")

    # ── S50: Choosing Operator ────────────────────────────
    build_diagram_slide(prs, 54, "Choosing the Right Operator",
        [("Where does the work live?", "Rule"),
         ("Shell command / script", "BashOperator"),
         ("Python function or logic", "PythonOperator"),
         ("SQL statement / database", "SQL operator (provider)"),
         ("REST API call", "HttpOperator or PythonOperator"),
         ("Hop pipeline via CLI", "BashOperator (calls hop-run)"),
         ("Docker container", "DockerOperator"),
         ("Wait for file", "FileSensor"),
         ("Wait for another DAG", "ExternalTaskSensor")],
        diagram_type="table",
        section_label="Operators & Sensors",
        note="If work is SQL → SQL. If Python → Python. If Hop CLI → Bash/Docker.")

    # ── S51: What is a Sensor ─────────────────────────────
    build_content_slide(prs, 55, "What is a Sensor?",
        ["A sensor is a special operator that waits for a condition before proceeding",
         "FileSensor: waits for a file to appear at a path",
         "HttpSensor: waits for an HTTP endpoint to return the expected response",
         "SqlSensor: waits for a SQL query to return a truthy result",
         "Sensors block downstream tasks from running until the condition is met"],
        section_label="Operators & Sensors",
        interactive_q="A sensor ______ for a condition.",
        interactive_a="waits")

    # ── S52: Poke vs Reschedule ───────────────────────────
    build_content_slide(prs, 56, "Sensor Mode: Poke vs Reschedule",
        ["poke mode: sensor holds a worker slot the entire time it waits",
         "reschedule mode: sensor releases its worker slot between checks",
         "Worker slots are like classroom seats — do not reserve one for hours while doing nothing",
         "Use reschedule for sensors that may wait hours or days",
         "Use poke only for very short waits (seconds to minutes)"],
        section_label="Operators & Sensors",
        interactive_q="Which sensor mode is more resource-efficient for long waits?",
        interactive_a="reschedule — it frees the worker slot between checks")

    # ── S53: Deferrable Operators ─────────────────────────
    build_content_slide(prs, 57, "Deferrable Operators",
        ["Modern way to wait asynchronously — requires Triggerer component",
         "Task suspends itself (defers) and releases worker slot",
         "Triggerer (lightweight async process) monitors the condition",
         "When condition met, Triggerer re-queues the task on a worker",
         "Most efficient for waits in production — minimal resource overhead"],
        section_label="Operators & Sensors",
        tip="Only mention deferrable operators after students understand poke/reschedule. It is the production-grade solution.")

    # ── S54: What is a Hook ───────────────────────────────
    build_content_slide(prs, 58, "What is a Hook?",
        ["A hook contains the connection logic to an external system",
         "Operator says WHAT to do; Hook knows HOW to connect",
         "Example: PostgresOperator uses PostgresHook to connect to the database",
         "Hooks reuse connection credentials stored as Airflow Connections",
         "Developers rarely write hooks directly — operators use them internally"],
        section_label="Operators & Sensors")

    # ── S55: Connections ──────────────────────────────────
    build_content_slide(prs, 59, "Connections — Securing System Access",
        ["Connections store system access info: host, port, login, password, schema",
         "Tasks reference a conn_id string — not hardcoded credentials",
         "Bad practice: password in DAG code (security risk, visible in logs)",
         "Good practice: conn_id='my_postgres' → credentials managed in UI or environment",
         "Store connections via UI (Admin → Connections) or environment variables"],
        section_label="Operators & Sensors",
        interactive_q="Why avoid hardcoded passwords in DAG code?",
        interactive_a="Security risk — visible in source control, logs, and anyone who reads the file")

    # ── S56: Variables ────────────────────────────────────
    build_content_slide(prs, 60, "Variables",
        ["Airflow Variables store small config values accessible in DAGs",
         "Examples: batch_size=1000, environment=prod, s3_bucket=my-data-lake",
         "Set via: UI (Admin → Variables), CLI, environment variables, or secrets backend",
         "Warning: Variable.get() at DAG top-level causes DB hit on every parse — use inside tasks",
         "For secrets (passwords, tokens), use Secrets backends instead"],
        section_label="Operators & Sensors")

    # ── S57: Secrets ──────────────────────────────────────
    build_content_slide(prs, 61, "Secrets Management",
        ["Sensitive values should never be stored in DAG code or plain Variables",
         "Airflow Secrets Backends: HashiCorp Vault, AWS Secrets Manager, GCP Secret Manager",
         "Secrets are retrieved at task runtime — not stored in Airflow DB",
         "For the lab: environment variable injection is sufficient",
         "Always separate credentials from code — this is non-negotiable in production"],
        section_label="Operators & Sensors",
        tip="Point to OWASP A2 (Cryptographic Failures) and A7 (Identification Failures) — hardcoded secrets violate both")

    # ── S58: XCom ─────────────────────────────────────────
    build_content_slide(prs, 62, "XCom — Cross-Task Communication",
        ["XCom = Cross-Communication — passes small data between tasks",
         "Sticky note analogy: task A writes a note; task B reads it",
         "Good for: row counts, file paths, status messages, small JSON",
         "Bad for: large DataFrames, files, or anything > a few KB",
         "XCom values are stored in the Metadata DB — not for big data"],
        section_label="Operators & Sensors",
        interactive_q="XCom should pass only ______ data.",
        interactive_a="small (metadata, paths, counts — not DataFrames or large files)")

    # ── S59: XCom Limitations ─────────────────────────────
    build_content_slide(prs, 63, "XCom Limitations",
        ["XCom values live in the Metadata DB — not designed for large payloads",
         "Large data should go to: S3/MinIO, database, file system, Kafka",
         "XCom can pass the file path or storage pointer — not the data itself",
         "XCom Backend can be extended to use S3/GCS for larger values",
         "Common mistake: pushing a full pandas DataFrame into XCom"],
        section_label="Operators & Sensors",
        tip="Rule: if it does not fit on a sticky note, do not put it in XCom")

    # ── S60: TaskFlow API ─────────────────────────────────
    build_lab_slide(prs, 64, "TaskFlow API — Pythonic DAGs",
        ["@dag decorator wraps the DAG definition function",
         "@task decorator turns a Python function into a task",
         "Return values automatically become XCom values",
         "Function arguments automatically pull XCom from upstream tasks",
         "Cleaner than traditional operators for Python-heavy workflows"],
        code_snippet="""from airflow.sdk import dag, task
import datetime

@dag(schedule="@daily",
     start_date=datetime.datetime(2024,1,1),
     catchup=False)
def my_pipeline():

    @task
    def extract():
        return {"rows": 100}

    @task
    def transform(data):
        return data["rows"] * 2

    @task
    def load(count):
        print(f"Loading {count} records")

    load(transform(extract()))

my_pipeline()""",
        section_label="Operators & Sensors",
        objective="Write clean Pythonic DAGs using decorators")

    # ── S61: Traditional vs TaskFlow ──────────────────────
    build_content_slide(prs, 65, "Traditional Operators vs TaskFlow API",
        ["Traditional: explicit operator classes (BashOperator, PythonOperator, etc.)",
         "TaskFlow: @task decorators for Python functions — cleaner, less boilerplate",
         "Both are valid and can be mixed in the same DAG",
         "Read both styles — you will encounter both in real codebases",
         "Prefer TaskFlow for new Python function workflows; use operators for non-Python work"],
        section_label="Operators & Sensors",
        interactive_q="Can you mix @task decorated functions with BashOperator in the same DAG?",
        interactive_a="Yes — traditional operators and TaskFlow can coexist in one DAG")

    # ── SECTION 5: Scheduling, Dependencies, Failure ──────
    build_section_divider(prs, 66, 5, "Scheduling, Dependencies & Failure",
        "When things run, why they might fail, and how to recover")

    # ── S62: What is Scheduling ───────────────────────────
    build_content_slide(prs, 67, "What is Scheduling?",
        ["Scheduling decides when a DAG run is created and eligible to execute",
         "Options: time-based (cron), manual trigger, event-driven (assets/datasets)",
         "Start with simple examples: @daily, @hourly, every Monday at 9 AM",
         "The Scheduler checks eligible DAGs every few seconds",
         "Schedule is not the same as execution start — tasks queue and execute after scheduling"],
        section_label="Scheduling & Failure",
        interactive_q="What is the difference between schedule and execution start?",
        interactive_a="Schedule creates the DAG run; actual task execution depends on available workers")

    # ── S63: Manual Trigger ───────────────────────────────
    build_content_slide(prs, 68, "Manual Trigger",
        ["A DAG can be triggered manually from the UI (Trigger DAG button)",
         "Also from CLI: airflow dags trigger <dag_id>",
         "Also from REST API: POST /api/v1/dags/{dag_id}/dagRuns",
         "Useful for: testing, emergency reruns, controlled backfills",
         "Manual run does not affect scheduled runs — they are independent"],
        section_label="Scheduling & Failure",
        tip="First thing to do in the lab: trigger manually. This removes scheduling confusion from early learning.")

    # ── S64: Cron Scheduling ──────────────────────────────
    build_content_slide(prs, 69, "Cron Scheduling",
        ["Cron expressions: minute  hour  day  month  day_of_week",
         "@daily = 0 0 * * *  (midnight every day)",
         "@hourly = 0 * * * *  (top of every hour)",
         "@weekly = 0 0 * * 0  (Sunday midnight)",
         "0 9 * * MON-FRI  (9 AM Monday to Friday)"],
        section_label="Scheduling & Failure",
        interactive_q="Decode: */15 * * * *",
        interactive_a="Every 15 minutes")

    # ── S65: Cron Examples ────────────────────────────────
    build_diagram_slide(prs, 70, "Cron Expressions — Reference Table",
        [("Expression", "Meaning"),
         ("None", "Manual trigger only"),
         ("@once", "Run once then never again"),
         ("@daily  /  0 0 * * *", "Every day at midnight"),
         ("@hourly  /  0 * * * *", "Every hour on the hour"),
         ("@weekly  /  0 0 * * 0", "Every Sunday at midnight"),
         ("0 9 * * MON-FRI", "Weekdays at 9:00 AM"),
         ("*/15 * * * *", "Every 15 minutes"),
         ("0 6 1 * *", "First day of month at 6 AM"),
         ("0 20 * * FRI", "Every Friday at 8 PM")],
        diagram_type="table",
        section_label="Scheduling & Failure")

    # ── S66: Logical Date ─────────────────────────────────
    build_content_slide(prs, 71, "Logical Date — Often Confusing!",
        ["Logical date (data_interval_start) identifies WHICH data interval is being processed",
         "It is NOT necessarily when you clicked run",
         "Report date analogy: you generate Monday report on Tuesday, but it belongs to Monday",
         "A daily DAG running at 01:00 on Jan 16 has logical_date = Jan 15 (yesterday's data)",
         "Access in code: {{ ds }} = logical date as YYYY-MM-DD string"],
        section_label="Scheduling & Failure",
        interactive_q="A daily DAG runs at midnight. Its logical_date is Jan 15. What data should it process?",
        interactive_a="Jan 15 data — the data interval for that specific run")

    # ── S67: Data Interval ────────────────────────────────
    build_content_slide(prs, 72, "Data Interval",
        ["Every scheduled run covers a specific time window of data",
         "Daily DAG: data_interval_start = 2024-01-15 00:00, data_interval_end = 2024-01-16 00:00",
         "Access: {{ data_interval_start }}, {{ data_interval_end }} in templates",
         "Use these to filter source data — ensures each run processes exactly its window",
         "Enables idempotent reruns — same interval, same result"],
        section_label="Scheduling & Failure")

    # ── S68: Assets Event Scheduling ──────────────────────
    build_content_slide(prs, 73, "Assets / Event-Driven Scheduling",
        ["Airflow can trigger a DAG when a data asset is produced — not just by clock",
         "Airflow 3 terminology: Assets; Airflow 2 terminology: Datasets (same concept)",
         "Producer DAG marks an asset as updated; Consumer DAG reacts automatically",
         "Enables data-driven pipelines without polling or complex sensor chains",
         "Example: bronze_load DAG updates bronze_orders asset → silver_transform starts"],
        section_label="Scheduling & Failure",
        tip="Ref: airflow.apache.org/docs/apache-airflow/stable/authoring-and-scheduling/assets.html")

    # ── S69: Trigger Rules ────────────────────────────────
    build_diagram_slide(prs, 74, "Trigger Rules",
        [("Trigger Rule", "When downstream task runs"),
         ("all_success (default)", "Only if ALL upstream tasks succeeded"),
         ("all_done", "After ALL upstream tasks finish (any state)"),
         ("all_failed", "Only if ALL upstream tasks failed"),
         ("one_success", "As soon as ONE upstream task succeeds"),
         ("one_failed", "As soon as ONE upstream task fails → use for alerts"),
         ("none_failed", "If no upstream task failed (skips are OK)"),
         ("none_skipped", "If no upstream task was skipped"),
         ("always", "Run regardless of upstream state")],
        diagram_type="table",
        section_label="Scheduling & Failure",
        note="Default all_success is strict — safest for ETL. none_failed useful when optional branches may be skipped.")

    # ── S70: Why Pipelines Fail ───────────────────────────
    build_content_slide(prs, 75, "Why Production Pipelines Fail",
        ["Network: timeout, DNS failure, connection refused",
         "Bad data: unexpected nulls, schema change, encoding issues",
         "Wrong credentials: expired token, rotated password",
         "Code bugs: index error, type mismatch, unhandled edge case",
         "Resource limits: memory OOM, disk full, worker slot exhausted",
         "External service: upstream DB down, API rate limit exceeded"],
        section_label="Scheduling & Failure",
        interactive_q="Normalise failure with your team: production systems fail. Good engineers ______ for failure.",
        interactive_a="plan (retries, alerts, runbooks, idempotency)")

    # ── S71: Retries ──────────────────────────────────────
    build_content_slide(prs, 76, "Retries — Recovering Automatically",
        ["retries=3 means Airflow retries the task 3 times before marking failed",
         "retry_delay=timedelta(minutes=5) waits 5 min between retries",
         "Good for: network timeouts, service temporarily unavailable, transient errors",
         "Bad for: permanent code bugs — retries waste time and increase noise",
         "Fix the code, not just the retry count"],
        section_label="Scheduling & Failure",
        interactive_q="A retry is useful for a temporary ______.",
        interactive_a="failure (transient error, timeout, brief service outage)")

    # ── S72: Exponential Backoff ──────────────────────────
    build_content_slide(prs, 77, "Exponential Backoff",
        ["Retry wait times increase after each failure",
         "Instead of retrying every 5 min (flooding a down service), wait 5, 10, 20 min",
         "Phone call analogy: if nobody answers, wait longer before calling again",
         "Set: retry_exponential_backoff=True in default_args or task",
         "Reduces thundering herd problem — many retries hitting a recovering service"],
        section_label="Scheduling & Failure")

    # ── S73: Timeouts ─────────────────────────────────────
    build_content_slide(prs, 78, "Timeouts",
        ["execution_timeout: max time a single task instance may run",
         "dagrun_timeout: max time a full DAG run may take",
         "If exceeded, Airflow marks the task/run as failed",
         "Every serious production task should have a reasonable timeout",
         "Prevents zombie tasks from blocking workers indefinitely"],
        section_label="Scheduling & Failure",
        tip="Ask: 'What if a data load runs forever? Without timeout, it blocks a worker slot until someone manually kills it'")

    # ── S74: Callbacks ────────────────────────────────────
    build_content_slide(prs, 79, "Callbacks — Alerting on Events",
        ["on_failure_callback: runs when a task fails (e.g., send Slack alert)",
         "on_success_callback: runs when a task succeeds",
         "on_retry_callback: runs before each retry",
         "on_skipped_callback: runs when task is skipped",
         "sla_miss_callback (DAG-level): runs when SLA deadline is missed"],
        section_label="Scheduling & Failure",
        interactive_q="Which callback would you use to send a PagerDuty alert on task failure?",
        interactive_a="on_failure_callback")

    # ── S75: Idempotency ──────────────────────────────────
    build_content_slide(prs, 80, "Idempotency — Safe to Retry",
        ["An idempotent task produces the same result when run multiple times",
         "If a task is retried, it should not create duplicate data",
         "Safe patterns: delete partition then reload, use MERGE/UPSERT",
         "Unsafe pattern: plain INSERT without dedup — retry creates duplicates",
         "Design all tasks for idempotency — retries and catchup require it"],
        section_label="Scheduling & Failure",
        interactive_q="Your load task does INSERT every run. It fails and retries. What goes wrong?",
        interactive_a="Duplicate rows — the same data is inserted twice. Use MERGE or delete-then-insert.")

    # ── S76: SLA and Freshness ────────────────────────────
    build_content_slide(prs, 81, "SLA and Data Freshness",
        ["SLA (Service Level Agreement): promise that data will be ready by a certain time",
         "Airflow sla parameter: how long after logical_date the task should complete",
         "sla_miss_callback fires if the task does not finish within the SLA",
         "Data teams care not just success/failure — also 'was the data ready on time?'",
         "Missing an SLA may trigger escalation even if the job eventually succeeds"],
        section_label="Scheduling & Failure")

    # ── S77: Catchup and Backfill ─────────────────────────
    build_content_slide(prs, 82, "Safe Backfill Strategy",
        ["Backfill = deliberately running a DAG for historical periods",
         "Command: airflow dags backfill -s 2024-01-01 -e 2024-01-31 my_dag",
         "Never backfill months without planning — it floods executor and external systems",
         "Limit concurrency: max_active_runs=1 prevents parallel historical runs",
         "Test one run first; validate data; then expand backfill range"],
        section_label="Scheduling & Failure",
        interactive_q="Before running a backfill for 6 months of data, what should you do first?",
        interactive_a="Test one historical run, validate the output, then set safe concurrency limits")

    # ── BREAK SLIDES ──────────────────────────────────────
    build_break_slide(prs, 83, "Morning Break", "15 Minutes", "10:30 — Back at 10:45")
    build_break_slide(prs, 84, "Lunch Break", "1 Hour", "13:00 — Back at 14:00")
    build_break_slide(prs, 85, "Afternoon Break", "30 Minutes", "15:30 — Back at 16:00")

    # ── SECTION 6: Kafka and Hop Integration ──────────────
    build_section_divider(prs, 86, 6, "Kafka & Hop Integration",
        "Streaming events, visual ETL, and Airflow as the orchestration layer")

    # ── S78: Why Integrate with Kafka ─────────────────────
    build_content_slide(prs, 87, "Why Integrate Airflow with Kafka?",
        ["Kafka handles high-throughput event streaming — Airflow cannot and should not",
         "Airflow orchestrates jobs AROUND Kafka, not inside it",
         "Pattern: Kafka streams to landing zone → Airflow validates readiness → runs transforms",
         "Airflow can: start a producer job, validate topic availability, schedule batch jobs on landed data",
         "Airflow should NOT: consume millions of events directly in a task"],
        section_label="Kafka & Hop",
        interactive_q="Best sentence: A) Airflow replaces Kafka  B) Kafka streams events and Airflow orchestrates workflows",
        interactive_a="B — each tool has a distinct role. They complement each other.")

    # ── S79: Kafka Refresher ──────────────────────────────
    build_content_slide(prs, 88, "Kafka Concepts Refresher",
        ["Topic: named log/stream of events (e.g., orders.cdc, inventory.updates)",
         "Producer: writes events to a topic",
         "Consumer: reads and processes events from a topic",
         "Partition: parallel lanes within a topic for scalability",
         "Consumer Group: set of consumers sharing work — each message processed once per group"],
        section_label="Kafka & Hop",
        tip="Use classroom mailbox analogy: topic = named mailbox, producer = person posting letters, consumer = person reading them")

    # ── S80: Airflow Kafka Patterns ───────────────────────
    build_content_slide(prs, 89, "Airflow + Kafka Orchestration Patterns",
        ["Pattern 1: Airflow triggers a Kafka producer script (BashOperator)",
         "Pattern 2: Airflow sensor checks if landing files from Kafka consumer exist",
         "Pattern 3: Airflow SQL check validates row count in Kafka-populated table",
         "Pattern 4: Airflow schedules batch transforms on Kafka-landed Bronze data",
         "Pattern 5: Airflow monitors DLQ row counts and alerts if bad events accumulate"],
        section_label="Kafka & Hop")

    # ── S81: Kafka to Bronze Architecture ─────────────────
    build_diagram_slide(prs, 90, "Kafka → Bronze → Airflow Orchestration",
        ["Source\nDatabase", "Debezium\nCDC Connector", "Kafka\nTopic",
         "Landing /\nBronze Storage", "Airflow\nOrchestration", "Hop /\nSpark Transforms",
         "Silver /\nGold Tables", "BI / Report\nLayer"],
        diagram_type="flow",
        section_label="Kafka & Hop",
        note="Ref: kafka.apache.org/documentation  |  Airflow orchestrates batch stages around the streaming pipeline")

    # ── S82: Schema Registry and DLQ ─────────────────────
    build_content_slide(prs, 91, "Schema Registry and Dead Letter Queue",
        ["Schema Registry: ensures producers and consumers agree on message format (Avro/JSON schema)",
         "Schema evolution: producers can add fields without breaking existing consumers",
         "Dead Letter Queue (DLQ): bad messages that failed processing go here for investigation",
         "Airflow integration: monitor DLQ row count; alert if it grows unexpectedly",
         "Keep Day 8 focus on orchestration — these are referenced, not deep-dived"],
        section_label="Kafka & Hop",
        interactive_q="What is the purpose of a Dead Letter Queue?",
        interactive_a="To capture bad/unprocessable messages for investigation — prevents data loss")

    # ── S83: What is Apache Hop ───────────────────────────
    build_content_slide(prs, 92, "What is Apache Hop?",
        ["Apache Hop is an open-source visual data integration platform",
         "Design ETL/ELT pipelines visually by dragging and connecting transforms",
         "Pipelines: data transformation flows (.hpl files)",
         "Workflows: control flow — run pipelines in sequence, handle errors (.hwf files)",
         "Runs standalone or triggered from CLI (hop-run) — how Airflow integrates"],
        section_label="Kafka & Hop",
        interactive_q="Airflow can trigger a Hop pipeline using the ______ command.",
        interactive_a="hop-run (BashOperator calling the hop-run CLI)")

    # ── S84: Hop Pipeline vs Workflow ─────────────────────
    build_content_slide(prs, 93, "Hop Pipeline vs Hop Workflow",
        ["Hop Pipeline (.hpl): data moves between transform steps (like a data flow)",
         "Hop Workflow (.hwf): controls WHEN pipelines run and handles branching/errors",
         "Pipeline example: read CSV → clean → filter → write to database",
         "Workflow example: run bronze pipeline → if success, run silver pipeline → else send alert",
         "Airflow can trigger either via hop-run CLI"],
        section_label="Kafka & Hop")

    # ── S85: hop-run ──────────────────────────────────────
    build_lab_slide(prs, 94, "hop-run — The CLI Integration Point",
        ["hop-run runs Hop pipelines and workflows from the command line",
         "Exit code 0 = success, non-zero = failure (Airflow uses this)",
         "BashOperator calls hop-run and monitors the exit code",
         "Parameters can be passed with -p param_name=value",
         "Full Hop project and environment paths must be specified"],
        code_snippet="""/opt/hop/hop-run.sh \\
  --project my_project \\
  --environment prod \\
  --file /opt/hop/pipelines/silver_transform.hpl \\
  -p processing_date={{ ds }} \\
  -p source_path=/data/bronze/orders""",
        section_label="Kafka & Hop",
        objective="Run a Hop pipeline from Airflow using BashOperator")

    # ── S86: Airflow Schedules Hop ────────────────────────
    build_content_slide(prs, 95, "Airflow Schedules Apache Hop",
        ["Airflow decides WHEN Hop runs: daily, after Bronze loads, on file arrival",
         "BashOperator executes hop-run CLI command inside the Airflow worker",
         "hop-run returns exit code → Airflow task success or failure",
         "On failure, Airflow retries as configured — Hop does not need retry logic",
         "Airflow logs capture all hop-run output for debugging"],
        section_label="Kafka & Hop")

    # ── S87: Passing Parameters to Hop ───────────────────
    build_content_slide(prs, 96, "Passing Parameters from Airflow to Hop",
        ["Use Jinja templates in BashOperator bash_command",
         "{{ ds }} = logical date as YYYY-MM-DD",
         "{{ data_interval_start }} = start of data interval",
         "{{ var.value.my_variable }} = Airflow Variable value",
         "Example: -p processing_date={{ ds }} passes the run date to the Hop pipeline"],
        section_label="Kafka & Hop",
        interactive_q="How would you pass today's date to a Hop pipeline from Airflow?",
        interactive_a="-p processing_date={{ ds }} in the BashOperator bash_command")

    # ── S88: Full Architecture ────────────────────────────
    build_diagram_slide(prs, 97, "Kafka + Hop + Airflow: Full Architecture",
        ["Source DB\n(Postgres)", "Debezium\nCDC", "Kafka\nTopic",
         "Bronze\nStorage", "Airflow\nDAG", "Hop Pipeline\n(Silver transform)",
         "Silver/Gold\nTables", "BI Report\nLayer"],
        diagram_type="flow",
        section_label="Kafka & Hop",
        note="Draw slowly: Source→Kafka→Bronze→Airflow triggers Hop→Warehouse→Report. Each tool has one responsibility.")

    # ── S89: Integration Recap ────────────────────────────
    build_content_slide(prs, 98, "Integration Recap: Tool Responsibilities",
        ["Kafka: streams events at high throughput — the highway",
         "Apache Hop: visually transforms data — the workshop",
         "Airflow: orchestrates when and how everything runs — the control tower",
         "Do NOT use one tool for every job",
         "Separation of concerns makes each tool easier to scale, debug, and replace"],
        section_label="Kafka & Hop",
        interactive_q="Repeat together: 'Kafka ______, Hop ______, Airflow ______.'",
        interactive_a="Kafka streams, Hop transforms, Airflow orchestrates")

    # ── SECTION 7: Hands-on Lab ───────────────────────────
    build_section_divider(prs, 99, 7, "Hands-on Lab",
        "Build and schedule a real pipeline from scratch")

    # ── S90: Lab Overview ─────────────────────────────────
    build_content_slide(prs, 100, "Lab Overview",
        ["Lab 1: Start Airflow with Docker Compose",
         "Lab 2: Create a Hello World DAG — trigger manually, inspect logs",
         "Lab 3: Build an ETL-style DAG with dependencies (extract → transform → load)",
         "Lab 4: Introduce intentional failure and observe retries",
         "Lab 5: Simulate Kafka landing check and Hop command pattern"],
        section_label="Lab",
        interactive_q="The Airflow UI usually runs locally on port ______.",
        interactive_a="8080")

    # ── S91: Prerequisites ────────────────────────────────
    build_content_slide(prs, 101, "Lab Prerequisites",
        ["Docker Desktop installed and running",
         "At least 4 GB RAM available for Docker",
         "Terminal (PowerShell / bash / WSL2)",
         "VS Code or any text editor for editing DAG files",
         "Run: docker --version and docker compose version to verify"],
        section_label="Lab",
        tip="Do this check before the session. Docker issues are the #1 lab blocker.")

    # ── S92: Start Airflow ────────────────────────────────
    build_lab_slide(prs, 102, "Lab 1: Start Airflow with Docker Compose",
        ["1. Create working directory: mkdir ~/airflow-day8",
         "2. Download official compose file from Airflow docs",
         "3. Create folders: dags/ logs/ plugins/ config/",
         "4. Set AIRFLOW_UID in .env file",
         "5. Run airflow-init then docker compose up -d",
         "6. Check health: docker compose ps",
         "7. Open http://localhost:8080 — login: airflow / airflow"],
        code_snippet="""mkdir -p ~/airflow-day8
cd ~/airflow-day8

curl -LfO 'https://airflow.apache.org/docs/apache-airflow/stable/docker-compose.yaml'

mkdir -p ./dags ./logs ./plugins ./config

echo "AIRFLOW_UID=$(id -u)" > .env

docker compose up airflow-init

docker compose up -d

docker compose ps""",
        section_label="Lab",
        objective="Run Airflow locally and open the Web UI")

    # ── S93: Hello World DAG ──────────────────────────────
    build_lab_slide(prs, 103, "Lab 2: Create a Hello World DAG",
        ["Create file: dags/day8_hello_world.py",
         "Define DAG with dag_id='day8_hello_world', schedule=None, catchup=False",
         "Add one BashOperator task: echo 'Hello from Airflow!'",
         "Wait ~30 seconds for DAG Processor to discover the file",
         "In UI: find the DAG, trigger it, open Graph → click task → view Logs"],
        code_snippet="""from datetime import datetime
from airflow import DAG
from airflow.providers.standard.operators.bash import BashOperator

with DAG(
    dag_id="day8_hello_world",
    start_date=datetime(2024, 1, 1),
    schedule=None,
    catchup=False,
    tags=["day8"],
) as dag:

    say_hello = BashOperator(
        task_id="say_hello",
        bash_command="echo 'Hello from Airflow Day 8!'",
    )""",
        section_label="Lab",
        objective="Prove Airflow can discover, schedule, and execute a task")

    # ── S94: ETL DAG ──────────────────────────────────────
    build_lab_slide(prs, 104, "Lab 3: ETL-Style DAG with Dependencies",
        ["Three tasks: extract → transform → load",
         "Add @daily schedule and catchup=False",
         "Set retries=2, retry_delay=1 minute in default_args",
         "Use BashOperator for all three tasks (shell simulates data work)",
         "Set dependencies: extract >> transform >> load"],
        code_snippet="""with DAG(
    dag_id="day8_simple_etl",
    start_date=datetime(2024, 1, 1),
    schedule="@daily",
    catchup=False,
    default_args={
        "retries": 2,
        "retry_delay": timedelta(minutes=1),
    },
) as dag:

    extract = BashOperator(
        task_id="extract",
        bash_command="echo 'Extracting...' && echo '100 rows' > /tmp/data.txt",
    )
    transform = BashOperator(
        task_id="transform",
        bash_command="cat /tmp/data.txt | tr a-z A-Z > /tmp/clean.txt",
    )
    load = BashOperator(
        task_id="load",
        bash_command="cat /tmp/clean.txt && echo 'Loaded!'",
    )
    extract >> transform >> load""",
        section_label="Lab",
        objective="Build and schedule an ETL pipeline with explicit dependencies")

    # ── S95: Failure and Retry Lab ────────────────────────
    build_lab_slide(prs, 105, "Lab 4: Intentional Failure and Retry",
        ["Temporarily change transform task to: bash_command='exit 1'",
         "Trigger the DAG manually and watch it fail",
         "Observe: task turns red, retry counter increments, delays between retries",
         "Check: task logs show the failure and each retry attempt",
         "Restore the original command and observe recovery"],
        code_snippet="""# Temporarily break transform:
transform = BashOperator(
    task_id="transform",
    bash_command="echo 'Simulating failure' && exit 1",
)

# Observe in UI:
# 1. Task turns orange (retrying)
# 2. After 2 retries, turns red (failed)
# 3. Load task does NOT run (blocked by upstream failure)
# 4. Check logs for each retry attempt

# Restore:
transform = BashOperator(
    task_id="transform",
    bash_command="cat /tmp/data.txt | tr a-z A-Z",
)""",
        section_label="Lab",
        objective="Experience retry behaviour and understand how failures propagate")

    # ── S96: Kafka and Hop Lab ────────────────────────────
    build_lab_slide(prs, 106, "Lab 5: Simulated Kafka & Hop Integration",
        ["Step 1: Add task that creates a simulated Kafka landing file",
         "Step 2: Add FileSensor (or BashOperator check) for the landing file",
         "Step 3: Add Hop command pattern (echo placeholder if Hop not installed)",
         "Step 4: Chain: produce_kafka_data >> check_landing >> run_hop_pipeline",
         "Step 5: Trigger and observe the full flow in the UI"],
        code_snippet="""# Simulate Kafka producer writing to landing zone
produce = BashOperator(
    task_id="simulate_kafka_landing",
    bash_command="mkdir -p /tmp/landing && echo '{\"orders\":100}' > /tmp/landing/orders.json",
)

# Check landing file exists
check = BashOperator(
    task_id="check_landing_file",
    bash_command="test -f /tmp/landing/orders.json && echo 'File ready!'",
)

# Simulate Hop pipeline (echo if Hop not installed)
hop = BashOperator(
    task_id="run_hop_silver_transform",
    bash_command="echo 'hop-run --file silver_transform.hpl -p date={{ ds }}'",
)

produce >> check >> hop""",
        section_label="Lab",
        objective="Simulate a Kafka → landing → Hop orchestration pattern")

    # ── S97: Lab Debugging ────────────────────────────────
    build_content_slide(prs, 107, "Lab Debugging Checklist",
        ["DAG not visible: check dags/ folder path and Python syntax errors",
         "Import error in UI: Admin → DAG Import Errors — shows the exact Python error",
         "Task stuck queued: check docker compose ps — is scheduler running?",
         "Task fails with no logs: check scheduler logs — docker compose logs scheduler",
         "Permission error on logs/: verify AIRFLOW_UID matches your user"],
        section_label="Lab",
        interactive_q="Your DAG file is saved but not visible in the UI. First two checks?",
        interactive_a="1. Check dags/ folder path is correct. 2. Check Admin → DAG Import Errors for Python syntax errors.")

    # ── SECTION 8: Summary ────────────────────────────────
    build_section_divider(prs, 108, 8, "Summary & Recap",
        "What you learned today and where to go next")

    # ── S98–S102: Recap Slides ────────────────────────────
    build_recap_slide(prs, 109, "Day Recap: What You Learned Today",
        [("Orchestration", "Airflow controls workflows — it does not replace Kafka, Hop, or databases"),
         ("Architecture", "Scheduler, DAG Processor, Metadata DB, Executor, Worker, UI — each has one job"),
         ("DAGs & Tasks", "Directed Acyclic Graphs: tasks, dependencies, no loops, visual graph"),
         ("Operators & Sensors", "Right operator for right work: Bash, Python, SQL, HTTP, sensors"),
         ("Scheduling & Failure", "Cron, retries, timeouts, callbacks, idempotency, SLA monitoring"),
         ("Integration", "Kafka streams, Hop transforms, Airflow orchestrates the full stack")])

    build_content_slide(prs, 110, "Day Recap: DAG Thinking",
        ["A DAG is a recipe — describe it in Python, execute it on schedule",
         "Tasks are stations; dependencies are the rails between them",
         "No cycles — every workflow must have a clear start and end",
         "Task states tell you exactly what happened: queued, running, success, failed, retry",
         "The Graph view in the UI is your operational dashboard"],
        section_label="Summary",
        interactive_q="Draw a DAG for: receive order → validate → process payment → send confirmation.",
        interactive_a="receive_order >> validate >> process_payment >> send_confirmation (linear ETL)")

    build_content_slide(prs, 111, "Day Recap: Operators and Connections",
        ["BashOperator for shell/CLI tools (including hop-run)",
         "PythonOperator / @task for Python logic",
         "SQL operators for warehouse queries",
         "Connections keep credentials out of code — always use conn_id",
         "Secrets backends for production credential management"],
        section_label="Summary",
        interactive_q="Classify: run dbt models, execute a shell script, call a REST API.",
        interactive_a="dbt=BashOperator/PythonOperator, shell=BashOperator, REST API=HttpOperator or PythonOperator")

    build_content_slide(prs, 112, "Day Recap: Scheduling and Failure Handling",
        ["A job without monitoring is not production-ready",
         "Retries handle transient failures — idempotency makes them safe",
         "Timeouts prevent zombie tasks from blocking workers",
         "Callbacks enable alerting on failure, retry, or SLA miss",
         "Catchup=False by default — plan backfills carefully with concurrency limits"],
        section_label="Summary")

    build_content_slide(prs, 113, "Day Recap: Kafka and Hop Integration",
        ["Tool triangle: Kafka streams → Hop transforms → Airflow orchestrates",
         "Airflow triggers Hop using BashOperator + hop-run CLI",
         "Pass processing dates with Jinja templates: {{ ds }}",
         "Monitor Kafka readiness via file/row count checks before transforms",
         "Keep each tool in its lane — this makes each easier to scale and debug"],
        section_label="Summary",
        interactive_q="Complete: Kafka _______, Hop _______, Airflow _______.",
        interactive_a="streams events, transforms data, orchestrates the workflow")

    # ── S103: Final Checklist ─────────────────────────────
    build_content_slide(prs, 114, "Exit Checklist — Can You Answer These?",
        ["What problem does Airflow solve that cron cannot?",
         "Name the 5 minimum components of Airflow and their roles",
         "What is a DAG? Why can it not have a cycle?",
         "What is the difference between an operator and a task instance?",
         "How do retries help? When are they NOT enough?",
         "Why is Kafka not replaced by Airflow?",
         "How does Airflow trigger a Hop pipeline?",
         "Where do you look first when a task fails?"],
        section_label="Summary",
        tip="Use this as an exit-ticket verbal quiz. Students should answer in one sentence each.")

    # ── S104–S108: Exercises ──────────────────────────────
    build_quiz_slide(prs, 115, "Quick Quiz 1: Architecture",
        [("Which component parses DAG files and stores to Metadata DB?",
          "A) Scheduler  B) DAG Processor  C) Worker  D) Triggerer",
          "B) DAG Processor"),
         ("Which component decides WHEN tasks should execute?",
          "A) Worker  B) Web UI  C) Scheduler  D) Executor",
          "C) Scheduler"),
         ("Where does Airflow store task run history?",
          "A) DAG files  B) Metadata Database  C) Kafka topic  D) Worker memory",
          "B) Metadata Database")],
        section_label="Summary")

    build_quiz_slide(prs, 116, "Quick Quiz 2: DAGs and Operators",
        [("DAG stands for Directed ______ Graph.",
          "A) Automatic  B) Acyclic  C) Async  D) Active",
          "B) Acyclic"),
         ("Which operator runs shell commands?",
          "A) PythonOperator  B) HttpOperator  C) BashOperator  D) EmailOperator",
          "C) BashOperator"),
         ("What does catchup=True do when start_date is 3 months ago?",
          "A) Nothing  B) May create many historical runs  C) Deletes old logs  D) Disables retries",
          "B) May create many historical runs")],
        section_label="Summary")

    build_quiz_slide(prs, 117, "Quick Quiz 3: Integration",
        [("Kafka is best described as:",
          "A) An orchestrator  B) An event streaming platform  C) A visual ETL tool  D) A scheduler",
          "B) An event streaming platform"),
         ("Which CLI tool does Airflow use to trigger a Hop pipeline?",
          "A) airflow dags run  B) hop-run  C) docker exec  D) kafka-topics",
          "B) hop-run"),
         ("Why should XCom NOT carry large DataFrames?",
          "A) XCom is only for strings  B) XCom values stored in Metadata DB — not for big data  C) XCom is deprecated  D) No reason",
          "B) XCom values stored in Metadata DB — not for big data")],
        section_label="Summary")

    # ── S109: Homework ────────────────────────────────────
    build_content_slide(prs, 118, "Homework",
        ["1. Create a DAG with 4 tasks: extract, validate, transform, load",
         "2. Add one intentional failure with 2 retries (configure retry_delay)",
         "3. Add one notification placeholder task with BashOperator (echo alert)",
         "4. Draw the Airflow + Kafka + Hop architecture in your own words",
         "5. Write 5 cron expressions: every day, every hour, Monday 9 AM, every 15 min, 1st of month",
         "6. Explain in 5 lines why Airflow is NOT a streaming engine"],
        section_label="Summary",
        tip="Review homework at start of Day 9. Ask students to explain one task to the class.")

    # ── S110: References ──────────────────────────────────
    build_content_slide(prs, 119, "Key References",
        ["Airflow architecture: airflow.apache.org/docs/apache-airflow/stable/core-concepts/overview.html",
         "DAGs: airflow.apache.org/docs/apache-airflow/stable/core-concepts/dags.html",
         "Tasks & operators: airflow.apache.org/docs/apache-airflow/stable/core-concepts/tasks.html",
         "TaskFlow API: airflow.apache.org/docs/apache-airflow/stable/tutorial/taskflow.html",
         "Connections: airflow.apache.org/docs/apache-airflow/stable/howto/connection.html",
         "Assets/event scheduling: airflow.apache.org/docs/apache-airflow/stable/authoring-and-scheduling/assets.html",
         "Docker Compose: airflow.apache.org/docs/apache-airflow/stable/howto/docker-compose/index.html",
         "Kafka docs: kafka.apache.org/documentation/",
         "hop-run: hop.apache.org/manual/latest/hop-run/index.html"],
        section_label="Summary")

    # ── S111: Looking Ahead ───────────────────────────────
    build_content_slide(prs, 120, "Looking Ahead — Day 9 Preview",
        ["Day 9: Advanced pipeline patterns and data quality",
         "dbt integration with Airflow — orchestrating SQL transformations",
         "Data quality frameworks: Great Expectations, dbt tests",
         "Dynamic DAGs: generating tasks programmatically from config",
         "Production monitoring: metrics, observability, alerting pipelines"],
        section_label="Summary",
        tip="End with: 'Today you learned the control room. Next session we wire in more sophisticated systems.'")

    # ──────────────────────────────────────────────────────
    return prs


if __name__ == "__main__":
    import sys
    output_path = r"C:\Temp\Day_08_Slides_Revamped.pptx"
    print("Generating revamped slides...")
    prs = generate_all_slides()
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
